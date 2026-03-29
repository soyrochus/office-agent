from __future__ import annotations

from copy import deepcopy
from dataclasses import dataclass
from pathlib import Path
from typing import Any

from offagent.adapters import pptx_adapter
from offagent.domain.locators import parse_locator, to_v2_locator
from offagent.domain.models import Capability, ChildSummary, ObjectPayload
from offagent.errors import InvalidArgumentsError, TargetNotFoundError

try:
    from pptx.opc.constants import RELATIONSHIP_TYPE as RT
except ModuleNotFoundError:  # pragma: no cover - exercised through dependency checks
    RT = None


@dataclass(frozen=True)
class _PptxTarget:
    canonical_locator: str
    object_type: str
    slide_number: int | None = None
    shape_id: int | None = None
    paragraph_index: int | None = None
    run_index: int | None = None
    row_index: int | None = None
    column_index: int | None = None


class PptxObjectResolver:
    def get_object(self, document_path: Path, locator: str) -> ObjectPayload:
        canonical = to_v2_locator(locator, file_type="pptx")
        presentation = pptx_adapter._open_presentation(document_path)
        target = _parse_pptx_target(canonical)

        if target.object_type == "presentation":
            return _build_presentation_payload(document_path, presentation)
        if target.object_type == "slide":
            return _build_slide_payload(document_path, presentation, target)
        if target.object_type == "notes":
            return _build_notes_payload(document_path, presentation, target)
        if target.object_type in {"shape", "text_shape", "image_shape", "table", "group_shape"}:
            return _build_shape_payload(document_path, presentation, target)
        if target.object_type == "paragraph":
            return _build_paragraph_payload(document_path, presentation, target)
        if target.object_type == "run":
            return _build_run_payload(document_path, presentation, target)
        if target.object_type == "table_row":
            return _build_table_row_payload(document_path, presentation, target)
        if target.object_type == "table_cell":
            return _build_table_cell_payload(document_path, presentation, target)
        raise InvalidArgumentsError(f"Unsupported PPTX locator: {locator}")

    def list_children(
        self,
        document_path: Path,
        locator: str,
        *,
        child_type: str | None = None,
        limit: int | None = None,
    ) -> list[ChildSummary]:
        canonical = to_v2_locator(locator, file_type="pptx")
        presentation = pptx_adapter._open_presentation(document_path)
        target = _parse_pptx_target(canonical)

        if target.object_type == "presentation":
            children = _presentation_children(presentation, child_type=child_type)
        elif target.object_type == "slide":
            children = _slide_children(presentation, target, child_type=child_type)
        elif target.object_type in {"shape", "text_shape"}:
            children = _text_shape_children(presentation, target, child_type=child_type)
        elif target.object_type == "paragraph":
            children = _paragraph_children(presentation, target, child_type=child_type)
        elif target.object_type == "table":
            children = _table_children(presentation, target, child_type=child_type)
        elif target.object_type == "table_row":
            children = _table_row_children(presentation, target, child_type=child_type)
        else:
            children = ()

        if limit is not None:
            return list(children[:limit])
        return list(children)

    def resolve_capabilities(self, document_path: Path, locator: str) -> frozenset[Capability]:
        del document_path
        canonical = to_v2_locator(locator, file_type="pptx")
        target = _parse_pptx_target(canonical)
        return _capabilities_for(target.object_type)


def add_slide(
    document_path: Path,
    *,
    layout_index: int | None = None,
    layout_name: str | None = None,
    output_path: Path,
) -> tuple[str, str, dict[str, Any]]:
    presentation = pptx_adapter._open_presentation(document_path)
    layout = _resolve_slide_layout(presentation, layout_index=layout_index, layout_name=layout_name)
    presentation.slides.add_slide(layout)
    slide_number = len(presentation.slides)
    presentation.save(output_path)
    return (
        f"pptx:slide:{slide_number}",
        f"Added slide {slide_number} using layout {layout.name!r}.",
        {"layout_name": layout.name, "layout_index": _layout_index(presentation, layout)},
    )


def duplicate_slide(
    document_path: Path,
    locator: str,
    *,
    position: int | None = None,
    output_path: Path,
) -> tuple[str, str, dict[str, Any]]:
    canonical = to_v2_locator(locator, file_type="pptx")
    target = _parse_pptx_target(canonical)
    if target.object_type != "slide":
        raise InvalidArgumentsError("pptx_duplicate_slide requires a slide locator.")

    presentation = pptx_adapter._open_presentation(document_path)
    copied_position = _duplicate_slide_in_presentation(presentation, target.slide_number, position)
    presentation.save(output_path)
    return (
        f"pptx:slide:{copied_position}",
        f"Duplicated {canonical} to slide {copied_position}.",
        {"source_locator": canonical, "position": copied_position},
    )


def set_slide_layout(
    document_path: Path,
    locator: str,
    *,
    layout_index: int | None = None,
    layout_name: str | None = None,
    output_path: Path,
) -> tuple[str, str, dict[str, Any]]:
    if RT is None:
        raise RuntimeError("python-pptx is required for PPTX operations.")

    canonical = to_v2_locator(locator, file_type="pptx")
    target = _parse_pptx_target(canonical)
    if target.object_type != "slide":
        raise InvalidArgumentsError("pptx_set_slide_layout requires a slide locator.")

    presentation = pptx_adapter._open_presentation(document_path)
    slide = _resolve_slide(presentation, target)
    layout = _resolve_slide_layout(presentation, layout_index=layout_index, layout_name=layout_name)
    old_layout_rids = [
        r_id
        for r_id, rel in slide.part.rels.items()
        if rel.reltype == RT.SLIDE_LAYOUT
    ]
    for r_id in old_layout_rids:
        slide.part.drop_rel(r_id)
    slide.part.relate_to(layout.part, RT.SLIDE_LAYOUT)
    presentation.save(output_path)
    return (
        canonical,
        f"Updated {canonical} to layout {layout.name!r}.",
        {"layout_name": layout.name, "layout_index": _layout_index(presentation, layout)},
    )


def add_text_shape(
    document_path: Path,
    locator: str,
    *,
    text: str,
    left: int,
    top: int,
    width: int,
    height: int,
    output_path: Path,
) -> tuple[str, str, dict[str, Any]]:
    canonical = to_v2_locator(locator, file_type="pptx")
    target = _parse_pptx_target(canonical)
    if target.object_type != "slide":
        raise InvalidArgumentsError("pptx_add_text_shape requires a slide locator.")

    presentation = pptx_adapter._open_presentation(document_path)
    slide = _resolve_slide(presentation, target)
    shape = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    shape.text_frame.text = text
    locator = _shape_locator(target.slide_number, shape, _shape_object_type(shape))
    presentation.save(output_path)
    return (
        locator,
        f"Added text shape {locator} to {canonical}.",
        {"text": text, "left": left, "top": top, "width": width, "height": height},
    )


def _build_presentation_payload(document_path: Path, presentation) -> ObjectPayload:
    slides = tuple(presentation.slides)
    preview = ""
    if slides:
        preview = _slide_preview(slides[0])
    return ObjectPayload(
        document=pptx_adapter._document_ref(document_path),
        locator="pptx:presentation",
        object_type="presentation",
        preview=preview,
        properties={"slide_count": len(slides)},
        capabilities=_capability_tuple("presentation"),
        child_summary=_presentation_children(presentation),
    )


def _build_slide_payload(document_path: Path, presentation, target: _PptxTarget) -> ObjectPayload:
    slide = _resolve_slide(presentation, target)
    bundle = pptx_adapter.get_slide_bundle(document_path, target.slide_number)
    layout_name = getattr(getattr(slide, "slide_layout", None), "name", None)
    return ObjectPayload(
        document=pptx_adapter._document_ref(document_path),
        locator=target.canonical_locator,
        object_type="slide",
        preview=bundle.preview,
        properties={
            "slide_number": target.slide_number,
            "layout_name": layout_name,
            "shape_count": len(slide.shapes),
            "text_block_count": len(bundle.text_blocks),
            "notes_text": bundle.notes_text,
        },
        capabilities=_capability_tuple("slide"),
        parent_locator="pptx:presentation",
        child_summary=_slide_children(presentation, target),
    )


def _build_notes_payload(document_path: Path, presentation, target: _PptxTarget) -> ObjectPayload:
    slide = _resolve_slide(presentation, target)
    notes_text = pptx_adapter._notes_text(slide)
    return ObjectPayload(
        document=pptx_adapter._document_ref(document_path),
        locator=target.canonical_locator,
        object_type="notes",
        preview=notes_text[:120],
        properties={"slide_number": target.slide_number, "text": notes_text},
        capabilities=_capability_tuple("notes"),
        parent_locator=f"pptx:slide:{target.slide_number}",
    )


def _build_shape_payload(document_path: Path, presentation, target: _PptxTarget) -> ObjectPayload:
    resolved = _resolve_shape_target(presentation, target)
    shape = resolved["shape"]
    actual_object_type = _shape_object_type(shape)
    if target.object_type != "shape" and target.object_type != actual_object_type:
        raise TargetNotFoundError(
            f"Shape {resolved['shape_id']} on slide {resolved['slide_number']} is not a {target.object_type}."
        )

    properties: dict[str, Any] = {
        "slide_number": resolved["slide_number"],
        "shape_id": resolved["shape_id"],
        "shape_index": resolved["shape_index"],
        "shape_name": getattr(shape, "name", None),
        "shape_type": getattr(getattr(shape, "shape_type", None), "name", str(getattr(shape, "shape_type", ""))),
        "left": int(getattr(shape, "left", 0)),
        "top": int(getattr(shape, "top", 0)),
        "width": int(getattr(shape, "width", 0)),
        "height": int(getattr(shape, "height", 0)),
        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
    }
    child_summary: tuple[ChildSummary, ...] = ()
    if getattr(shape, "has_text_frame", False):
        properties["text"] = pptx_adapter._text_frame_text(shape.text_frame)
    if getattr(shape, "has_table", False):
        table = shape.table
        properties["row_count"] = len(table.rows)
        properties["column_count"] = len(table.columns)
        child_summary = _table_children(presentation, _PptxTarget(target.canonical_locator, "table", target.slide_number, target.shape_id))

    return ObjectPayload(
        document=pptx_adapter._document_ref(document_path),
        locator=target.canonical_locator,
        object_type=actual_object_type if target.object_type == "shape" else target.object_type,
        preview=_shape_preview(shape),
        properties=properties,
        capabilities=_capability_tuple(actual_object_type if target.object_type == "shape" else target.object_type),
        parent_locator=f"pptx:slide:{resolved['slide_number']}",
        child_summary=child_summary,
    )


def _build_paragraph_payload(document_path: Path, presentation, target: _PptxTarget) -> ObjectPayload:
    resolved = _resolve_paragraph_target(presentation, target)
    paragraph = resolved["paragraph"]
    return ObjectPayload(
        document=pptx_adapter._document_ref(document_path),
        locator=target.canonical_locator,
        object_type="paragraph",
        preview=paragraph.text[:120],
        properties={
            "slide_number": resolved["slide_number"],
            "shape_id": resolved["shape_id"],
            "paragraph_index": resolved["paragraph_index"],
            "text": paragraph.text,
            "run_count": len(paragraph.runs),
        },
        capabilities=_capability_tuple("paragraph"),
        parent_locator=f"pptx:slide:{resolved['slide_number']}:text_shape:{resolved['shape_id']}",
        child_summary=_paragraph_children(presentation, target),
    )


def _build_run_payload(document_path: Path, presentation, target: _PptxTarget) -> ObjectPayload:
    resolved = _resolve_run_target(presentation, target)
    run = resolved["run"]
    return ObjectPayload(
        document=pptx_adapter._document_ref(document_path),
        locator=target.canonical_locator,
        object_type="run",
        preview=run.text[:120],
        properties={
            "slide_number": resolved["slide_number"],
            "shape_id": resolved["shape_id"],
            "paragraph_index": resolved["paragraph_index"],
            "run_index": resolved["run_index"],
            "text": run.text,
        },
        capabilities=_capability_tuple("run"),
        parent_locator=f"pptx:slide:{resolved['slide_number']}:text_shape:{resolved['shape_id']}:para:{resolved['paragraph_index']}",
    )


def _build_table_row_payload(document_path: Path, presentation, target: _PptxTarget) -> ObjectPayload:
    resolved = _resolve_table_row_target(presentation, target)
    return ObjectPayload(
        document=pptx_adapter._document_ref(document_path),
        locator=target.canonical_locator,
        object_type="table_row",
        preview=" | ".join(cell.text for cell in resolved["row"].cells)[:120],
        properties={
            "slide_number": resolved["slide_number"],
            "shape_id": resolved["shape_id"],
            "row_index": resolved["row_index"],
            "cell_count": len(resolved["row"].cells),
        },
        capabilities=_capability_tuple("table_row"),
        parent_locator=f"pptx:slide:{resolved['slide_number']}:table:{resolved['shape_id']}",
        child_summary=_table_row_children(presentation, target),
    )


def _build_table_cell_payload(document_path: Path, presentation, target: _PptxTarget) -> ObjectPayload:
    resolved = _resolve_table_cell_target(presentation, target)
    cell = resolved["cell"]
    return ObjectPayload(
        document=pptx_adapter._document_ref(document_path),
        locator=target.canonical_locator,
        object_type="table_cell",
        preview=cell.text[:120],
        properties={
            "slide_number": resolved["slide_number"],
            "shape_id": resolved["shape_id"],
            "row_index": resolved["row_index"],
            "column_index": resolved["column_index"],
            "text": cell.text,
        },
        capabilities=_capability_tuple("table_cell"),
        parent_locator=f"pptx:slide:{resolved['slide_number']}:table:{resolved['shape_id']}:row:{resolved['row_index']}",
    )


def _presentation_children(presentation, *, child_type: str | None = None) -> tuple[ChildSummary, ...]:
    if child_type not in {None, "", "slide"}:
        return ()
    return tuple(
        ChildSummary(
            locator=f"pptx:slide:{slide_number}",
            object_type="slide",
            preview=_slide_preview(slide),
            capabilities=_capability_tuple("slide"),
        )
        for slide_number, slide in enumerate(presentation.slides, start=1)
    )


def _slide_children(presentation, target: _PptxTarget, *, child_type: str | None = None) -> tuple[ChildSummary, ...]:
    slide = _resolve_slide(presentation, target)
    normalized_child_type = child_type or None
    children: list[ChildSummary] = []

    if normalized_child_type in {None, "notes"}:
        children.append(
            ChildSummary(
                locator=f"pptx:slide:{target.slide_number}:notes",
                object_type="notes",
                preview=pptx_adapter._notes_text(slide)[:120],
                capabilities=_capability_tuple("notes"),
            )
        )

    for shape in slide.shapes:
        object_type = _shape_object_type(shape)
        if normalized_child_type not in {None, "shape", object_type}:
            continue
        children.append(
            ChildSummary(
                locator=_shape_locator(target.slide_number, shape, object_type),
                object_type=object_type,
                preview=_shape_preview(shape),
                capabilities=_capability_tuple(object_type),
            )
        )
    return tuple(children)


def _text_shape_children(presentation, target: _PptxTarget, *, child_type: str | None = None) -> tuple[ChildSummary, ...]:
    if child_type not in {None, "", "paragraph"}:
        return ()
    resolved = _resolve_shape_target(presentation, target)
    shape = resolved["shape"]
    if not getattr(shape, "has_text_frame", False):
        return ()
    return tuple(
        ChildSummary(
            locator=f"pptx:slide:{resolved['slide_number']}:text_shape:{resolved['shape_id']}:para:{paragraph_index}",
            object_type="paragraph",
            preview=paragraph.text[:120],
            capabilities=_capability_tuple("paragraph"),
        )
        for paragraph_index, paragraph in enumerate(shape.text_frame.paragraphs)
    )


def _paragraph_children(presentation, target: _PptxTarget, *, child_type: str | None = None) -> tuple[ChildSummary, ...]:
    if child_type not in {None, "", "run"}:
        return ()
    resolved = _resolve_paragraph_target(presentation, target)
    return tuple(
        ChildSummary(
            locator=(
                f"pptx:slide:{resolved['slide_number']}:text_shape:{resolved['shape_id']}:para:{resolved['paragraph_index']}:run:{run_index}"
            ),
            object_type="run",
            preview=run.text[:120],
            capabilities=_capability_tuple("run"),
        )
        for run_index, run in enumerate(resolved["paragraph"].runs)
    )


def _table_children(presentation, target: _PptxTarget, *, child_type: str | None = None) -> tuple[ChildSummary, ...]:
    resolved = _resolve_shape_target(presentation, target)
    shape = resolved["shape"]
    if not getattr(shape, "has_table", False) or child_type not in {None, "", "table_row"}:
        return ()
    return tuple(
        ChildSummary(
            locator=f"pptx:slide:{resolved['slide_number']}:table:{resolved['shape_id']}:row:{row_index}",
            object_type="table_row",
            preview=" | ".join(cell.text for cell in row.cells)[:120],
            capabilities=_capability_tuple("table_row"),
        )
        for row_index, row in enumerate(shape.table.rows)
    )


def _table_row_children(presentation, target: _PptxTarget, *, child_type: str | None = None) -> tuple[ChildSummary, ...]:
    resolved = _resolve_table_row_target(presentation, target)
    if child_type not in {None, "", "table_cell"}:
        return ()
    return tuple(
        ChildSummary(
            locator=(
                f"pptx:slide:{resolved['slide_number']}:table:{resolved['shape_id']}:row:{resolved['row_index']}:cell:{column_index}"
            ),
            object_type="table_cell",
            preview=cell.text[:120],
            capabilities=_capability_tuple("table_cell"),
        )
        for column_index, cell in enumerate(resolved["row"].cells)
    )


def _resolve_slide(presentation, target: _PptxTarget):
    assert target.slide_number is not None
    return pptx_adapter._resolve_slide(presentation, target.slide_number)


def _resolve_shape_target(presentation, target: _PptxTarget) -> dict[str, Any]:
    slide = _resolve_slide(presentation, target)
    assert target.shape_id is not None
    for shape_index, shape in enumerate(slide.shapes):
        if shape.shape_id == target.shape_id:
            return {
                "slide_number": target.slide_number,
                "shape_id": target.shape_id,
                "shape_index": shape_index,
                "shape": shape,
            }
    raise TargetNotFoundError(f"Shape {target.shape_id} does not exist on slide {target.slide_number}.")


def _resolve_paragraph_target(presentation, target: _PptxTarget) -> dict[str, Any]:
    resolved = _resolve_shape_target(presentation, target)
    shape = resolved["shape"]
    if not getattr(shape, "has_text_frame", False):
        raise TargetNotFoundError(f"Shape {resolved['shape_id']} is not a text shape.")
    assert target.paragraph_index is not None
    try:
        paragraph = shape.text_frame.paragraphs[target.paragraph_index]
    except IndexError as exc:
        raise TargetNotFoundError(
            f"Paragraph {target.paragraph_index} does not exist in shape {resolved['shape_id']}."
        ) from exc
    return {**resolved, "paragraph_index": target.paragraph_index, "paragraph": paragraph}


def _resolve_run_target(presentation, target: _PptxTarget) -> dict[str, Any]:
    resolved = _resolve_paragraph_target(presentation, target)
    assert target.run_index is not None
    try:
        run = resolved["paragraph"].runs[target.run_index]
    except IndexError as exc:
        raise TargetNotFoundError(
            f"Run {target.run_index} does not exist in paragraph {resolved['paragraph_index']}."
        ) from exc
    return {**resolved, "run_index": target.run_index, "run": run}


def _resolve_table_row_target(presentation, target: _PptxTarget) -> dict[str, Any]:
    resolved = _resolve_shape_target(presentation, target)
    shape = resolved["shape"]
    if not getattr(shape, "has_table", False):
        raise TargetNotFoundError(f"Shape {resolved['shape_id']} is not a table.")
    assert target.row_index is not None
    try:
        row = shape.table.rows[target.row_index]
    except IndexError as exc:
        raise TargetNotFoundError(
            f"Row {target.row_index} does not exist in table {resolved['shape_id']}."
        ) from exc
    return {**resolved, "row_index": target.row_index, "row": row}


def _resolve_table_cell_target(presentation, target: _PptxTarget) -> dict[str, Any]:
    resolved = _resolve_table_row_target(presentation, target)
    assert target.column_index is not None
    try:
        cell = resolved["row"].cells[target.column_index]
    except IndexError as exc:
        raise TargetNotFoundError(
            f"Cell {target.column_index} does not exist in row {resolved['row_index']}."
        ) from exc
    return {**resolved, "column_index": target.column_index, "cell": cell}


def _shape_object_type(shape) -> str:
    if getattr(shape, "has_table", False):
        return "table"
    shape_type_name = getattr(getattr(shape, "shape_type", None), "name", "")
    if shape_type_name == "PICTURE":
        return "image_shape"
    if shape_type_name == "GROUP":
        return "group_shape"
    if getattr(shape, "has_text_frame", False):
        return "text_shape"
    return "shape"


def _shape_locator(slide_number: int, shape, object_type: str) -> str:
    if object_type == "shape":
        return f"pptx:slide:{slide_number}:shape:{shape.shape_id}"
    return f"pptx:slide:{slide_number}:{object_type}:{shape.shape_id}"


def _shape_preview(shape) -> str:
    if getattr(shape, "has_text_frame", False):
        return pptx_adapter._text_frame_text(shape.text_frame)[:120]
    if getattr(shape, "has_table", False):
        rows = shape.table.rows
        if rows:
            first_row = rows[0]
            return " | ".join(cell.text for cell in first_row.cells)[:120]
    return getattr(shape, "name", "")[:120]


def _slide_preview(slide) -> str:
    text_blocks = pptx_adapter._slide_text_blocks(slide)
    return next((block.text[:120] for block in text_blocks if block.text), "")


def _resolve_slide_layout(presentation, *, layout_index: int | None, layout_name: str | None):
    layouts = tuple(presentation.slide_layouts)
    if (layout_index is None) == (layout_name is None):
        raise InvalidArgumentsError("Specify exactly one of layout_index or layout_name.")

    if layout_index is not None:
        if layout_index < 0 or layout_index >= len(layouts):
            raise InvalidArgumentsError(f"Unknown PPTX slide layout index: {layout_index}")
        return layouts[layout_index]

    assert layout_name is not None
    for layout in layouts:
        if layout.name == layout_name:
            return layout
    raise InvalidArgumentsError(f"Unknown PPTX slide layout: {layout_name}")


def _layout_index(presentation, layout) -> int:
    for index, candidate in enumerate(presentation.slide_layouts):
        if candidate == layout:
            return index
    raise RuntimeError("Failed to resolve PPTX slide layout index.")


def _duplicate_slide_in_presentation(presentation, slide_number: int | None, position: int | None) -> int:
    if slide_number is None:
        raise InvalidArgumentsError("pptx_duplicate_slide requires a slide locator.")

    source_slide = pptx_adapter._resolve_slide(presentation, slide_number)
    new_slide = presentation.slides.add_slide(source_slide.slide_layout)

    for placeholder_shape in list(new_slide.shapes):
        placeholder_shape.element.getparent().remove(placeholder_shape.element)

    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(deepcopy(shape.element), "p:extLst")

    for rel in source_slide.part.rels.values():
        if rel.reltype.endswith("/notesSlide") or rel.reltype.endswith("/slideLayout"):
            continue
        if rel.is_external:
            new_rid = new_slide.part.relate_to(rel.target_ref, rel.reltype, is_external=True)
        else:
            new_rid = new_slide.part.relate_to(rel.target_part, rel.reltype)
        _retarget_shape_relationships(new_slide, rel.rId, new_rid)

    if getattr(source_slide, "notes_slide", None) is not None:
        source_notes = getattr(source_slide.notes_slide, "notes_text_frame", None)
        target_notes = getattr(new_slide.notes_slide, "notes_text_frame", None)
        if source_notes is not None and target_notes is not None:
            target_notes.text = source_notes.text

    copied_position = len(presentation.slides) if position is None else position
    _move_slide_in_memory(presentation, len(presentation.slides), copied_position)
    return copied_position


def _move_slide_in_memory(presentation, slide_number: int, new_position: int) -> None:
    slide_count = len(presentation.slides)
    if new_position < 1 or new_position > slide_count:
        raise InvalidArgumentsError(f"Invalid target slide position: {new_position}")
    sld_id_list = presentation.slides._sldIdLst
    slide_id = sld_id_list.sldId_lst[slide_number - 1]
    sld_id_list.remove(slide_id)
    sld_id_list.insert(new_position - 1, slide_id)


def _retarget_shape_relationships(slide, source_rid: str, target_rid: str) -> None:
    for shape in slide.shapes:
        for element in shape.element.iter():
            for attr_name, attr_value in list(element.attrib.items()):
                if attr_value == source_rid:
                    element.set(attr_name, target_rid)


def _parse_pptx_target(locator: str) -> _PptxTarget:
    parsed = parse_locator(locator)
    components = parsed.components
    if components == ("pptx", "presentation"):
        return _PptxTarget(locator, "presentation")
    if len(components) == 3 and components[:2] == ("pptx", "slide"):
        return _PptxTarget(locator, "slide", slide_number=_require_index(components[2], locator))
    if len(components) == 4 and components[:2] == ("pptx", "slide") and components[3] == "notes":
        return _PptxTarget(locator, "notes", slide_number=_require_index(components[2], locator))
    if len(components) == 5 and components[:2] == ("pptx", "slide"):
        return _PptxTarget(
            locator,
            components[3],
            slide_number=_require_index(components[2], locator),
            shape_id=_require_index(components[4], locator),
        )
    if (
        len(components) == 7
        and components[:2] == ("pptx", "slide")
        and components[3] in {"shape", "text_shape"}
        and components[5] == "para"
    ):
        return _PptxTarget(
            locator,
            "paragraph",
            slide_number=_require_index(components[2], locator),
            shape_id=_require_index(components[4], locator),
            paragraph_index=_require_index(components[6], locator),
        )
    if (
        len(components) == 9
        and components[:2] == ("pptx", "slide")
        and components[3] in {"shape", "text_shape"}
        and components[5] == "para"
        and components[7] == "run"
    ):
        return _PptxTarget(
            locator,
            "run",
            slide_number=_require_index(components[2], locator),
            shape_id=_require_index(components[4], locator),
            paragraph_index=_require_index(components[6], locator),
            run_index=_require_index(components[8], locator),
        )
    if len(components) == 7 and components[:2] == ("pptx", "slide") and components[3] == "table" and components[5] == "row":
        return _PptxTarget(
            locator,
            "table_row",
            slide_number=_require_index(components[2], locator),
            shape_id=_require_index(components[4], locator),
            row_index=_require_index(components[6], locator),
        )
    if (
        len(components) == 9
        and components[:2] == ("pptx", "slide")
        and components[3] == "table"
        and components[5] == "row"
        and components[7] == "cell"
    ):
        return _PptxTarget(
            locator,
            "table_cell",
            slide_number=_require_index(components[2], locator),
            shape_id=_require_index(components[4], locator),
            row_index=_require_index(components[6], locator),
            column_index=_require_index(components[8], locator),
        )
    raise InvalidArgumentsError(f"Unsupported PPTX locator: {locator}")


def _capabilities_for(object_type: str) -> frozenset[Capability]:
    if object_type == "presentation":
        return frozenset({Capability.READ, Capability.ADD_CHILD})
    if object_type == "slide":
        return frozenset(
            {
                Capability.READ,
                Capability.UPDATE,
                Capability.DELETE,
                Capability.ADD_CHILD,
                Capability.MOVE,
                Capability.COPY,
            }
        )
    if object_type == "notes":
        return frozenset({Capability.READ, Capability.UPDATE})
    if object_type in {"shape", "text_shape", "image_shape", "group_shape"}:
        return frozenset({Capability.READ, Capability.UPDATE, Capability.DELETE, Capability.MOVE, Capability.COPY})
    if object_type == "paragraph":
        return frozenset({Capability.READ, Capability.UPDATE, Capability.STYLE})
    if object_type == "run":
        return frozenset({Capability.READ, Capability.UPDATE, Capability.STYLE})
    if object_type == "table":
        return frozenset(
            {
                Capability.READ,
                Capability.UPDATE,
                Capability.DELETE,
                Capability.ADD_CHILD,
                Capability.MOVE,
                Capability.COPY,
            }
        )
    if object_type == "table_row":
        return frozenset({Capability.READ, Capability.UPDATE, Capability.DELETE, Capability.MOVE, Capability.COPY})
    if object_type == "table_cell":
        return frozenset({Capability.READ, Capability.UPDATE, Capability.STYLE})
    return frozenset({Capability.READ})


def _capability_tuple(object_type: str) -> tuple[Capability, ...]:
    return tuple(sorted(_capabilities_for(object_type), key=lambda capability: capability.value))


def _require_index(raw: str, locator: str) -> int:
    try:
        return int(raw)
    except ValueError as exc:
        raise InvalidArgumentsError(f"Invalid PPTX locator: {locator}") from exc
