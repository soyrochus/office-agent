from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path

from offagent.domain.locators import parse_locator, to_v2_locator
from offagent.domain.models import (
    BlockStyle,
    DocumentRef,
    InlineStyle,
    IndexedItem,
    PresentationSlideSummary,
    PptxTextBlockNode,
    SectionPayload,
    SlideBundle,
    SlideTextBlock,
    StructureSection,
)
from offagent.errors import InvalidArgumentsError, TargetNotEditableError as BaseTargetNotEditableError
from offagent.errors import TargetNotFoundError

try:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
    from pptx.util import Pt
except ModuleNotFoundError:  # pragma: no cover - exercised through dependency checks
    Presentation = None
    RGBColor = None
    MSO_ANCHOR = None
    PP_ALIGN = None
    Pt = None


@dataclass(frozen=True)
class ResolvedShape:
    slide_number: int
    shape_id: int
    shape_index: int
    shape_name: str | None
    is_placeholder: bool
    text: str


class TargetNotEditableError(BaseTargetNotEditableError):
    """Raised when a requested PPTX target exists but is not a text frame."""


def extract_document(document_path: Path) -> list[IndexedItem]:
    presentation = _open_presentation(document_path)
    items: list[IndexedItem] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        for shape_index, shape in enumerate(slide.shapes):
            if not getattr(shape, "has_text_frame", False):
                continue

            item_id = make_item_id(slide_number, shape.shape_id)
            text = _text_frame_text(shape.text_frame)
            items.append(
                IndexedItem(
                    item_id=item_id,
                    item_type="slide_text_shape",
                    locator=item_id,
                    preview=text[:120],
                    content_text=text,
                    metadata={
                        "slide_number": slide_number,
                        "shape_id": shape.shape_id,
                        "shape_index": shape_index,
                        "shape_name": getattr(shape, "name", None),
                        "text_frame_text": text,
                        "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
                    },
                )
            )

    return items


def build_embedding_text(item: IndexedItem, document_path: Path) -> str:
    del document_path
    return item.content_text


def read_text_shape(document_path: Path, item_id: str) -> str:
    resolved = resolve_shape(document_path, item_id)
    return resolved.text


def replace_text_shape(document_path: Path, item_id: str, text: str, output_path: Path | None = None) -> Path:
    presentation = _open_presentation(document_path)
    shape = _resolve_shape(presentation, item_id)
    text_frame = _require_text_frame(shape)
    text_frame.clear()
    text_frame.paragraphs[0].text = text
    target_path = _target_path(document_path, output_path)
    presentation.save(target_path)
    return target_path


def append_text_shape(document_path: Path, item_id: str, text: str, output_path: Path | None = None) -> Path:
    presentation = _open_presentation(document_path)
    shape = _resolve_shape(presentation, item_id)
    text_frame = _require_text_frame(shape)
    text_frame.text = f"{_text_frame_text(text_frame)}{text}"
    target_path = _target_path(document_path, output_path)
    presentation.save(target_path)
    return target_path


def make_slide_locator(slide_number: int) -> str:
    return f"slide:{slide_number}"


def resolve_shape(document_path: Path, item_id: str) -> ResolvedShape:
    presentation = _open_presentation(document_path)
    shape = _resolve_shape(presentation, item_id)
    text_frame = _require_text_frame(shape)
    slide_number, shape_id = parse_item_id(item_id)
    return ResolvedShape(
        slide_number=slide_number,
        shape_id=shape_id,
        shape_index=_shape_index(shape),
        shape_name=getattr(shape, "name", None),
        is_placeholder=bool(getattr(shape, "is_placeholder", False)),
        text=_text_frame_text(text_frame),
    )


def resolve_structure(document_path: Path) -> tuple[StructureSection, ...]:
    presentation = _open_presentation(document_path)
    sections: list[StructureSection] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_blocks = _slide_text_blocks(slide)
        preview = next((block.text for block in text_blocks if block.text), "")
        locator = (
            make_item_id(slide_number, text_blocks[0].shape_id)
            if text_blocks
            else make_slide_locator(slide_number)
        )
        sections.append(
            StructureSection(
                locator=locator,
                section_type="slide",
                preview=preview[:120],
                metadata={
                    "slide_number": slide_number,
                    "shape_count": len(slide.shapes),
                    "text_block_count": len(text_blocks),
                },
            )
        )

    return tuple(sections)


def get_section(document_path: Path, locator: str) -> SectionPayload:
    slide_number = _slide_number_from_locator(locator)
    bundle = get_slide_bundle(document_path, slide_number)
    return SectionPayload(
        document=bundle.document,
        locator=locator if locator.startswith("slide:") else make_slide_locator(slide_number),
        section_type="slide",
        preview=bundle.preview,
        metadata=bundle.metadata,
        slide_number=bundle.slide_number,
        notes_text=bundle.notes_text,
        text_blocks=tuple(
            PptxTextBlockNode(
                locator=make_item_id(slide_number, block.shape_id),
                position=block.position,
                shape_id=block.shape_id,
                shape_name=block.shape_name,
                preview=block.preview,
                text=block.text,
                metadata=block.metadata,
            )
            for block in bundle.text_blocks
        ),
    )


def read_node(document_path: Path, locator: str) -> tuple[str, str, dict[str, object]]:
    normalized = locator.strip()
    if normalized.startswith("slide:") and ":shape:" not in normalized:
        slide_number = _slide_number_from_locator(normalized)
        bundle = get_slide_bundle(document_path, slide_number)
        text = "\n\n".join(block.text for block in bundle.text_blocks if block.text)
        return (
            "slide",
            text,
            {
                "slide_number": slide_number,
                "notes_text": bundle.notes_text,
                "text_block_count": len(bundle.text_blocks),
            },
        )

    resolved = resolve_shape(document_path, normalized)
    return (
        "slide_text_shape",
        resolved.text,
        {
            "slide_number": resolved.slide_number,
            "shape_id": resolved.shape_id,
            "shape_index": resolved.shape_index,
            "shape_name": resolved.shape_name,
            "is_placeholder": resolved.is_placeholder,
        },
    )


def write_node(document_path: Path, locator: str, text: str, output_path: Path | None = None) -> Path:
    normalized = locator.strip()
    if normalized.startswith("slide:") and ":shape:" not in normalized:
        shape_locator = _first_text_shape_locator(document_path, _slide_number_from_locator(normalized))
        if shape_locator is None:
            raise TargetNotEditableError("slide has no editable text shapes")
        return replace_text_shape(document_path, shape_locator, text, output_path)
    return replace_text_shape(document_path, normalized, text, output_path)


def get_presentation_structure(document_path: Path) -> tuple[PresentationSlideSummary, ...]:
    presentation = _open_presentation(document_path)
    slides: list[PresentationSlideSummary] = []

    for slide_number, slide in enumerate(presentation.slides, start=1):
        text_blocks = _slide_text_blocks(slide)
        preview = next((block.text for block in text_blocks if block.text), "")
        slides.append(
            PresentationSlideSummary(
                slide_number=slide_number,
                preview=preview[:120],
                metadata={
                    "slide_number": slide_number,
                    "shape_count": len(slide.shapes),
                    "text_block_count": len(text_blocks),
                },
            )
        )

    return tuple(slides)


def get_slide_bundle(document_path: Path, slide_number: int) -> SlideBundle:
    presentation = _open_presentation(document_path)
    slide = _resolve_slide(presentation, slide_number)
    text_blocks = _slide_text_blocks(slide)
    preview = next((block.text for block in text_blocks if block.text), "")
    return SlideBundle(
        document=_document_ref(document_path),
        slide_number=slide_number,
        preview=preview[:120],
        notes_text=_notes_text(slide),
        metadata={
            "slide_number": slide_number,
            "shape_count": len(slide.shapes),
            "text_block_count": len(text_blocks),
        },
        text_blocks=tuple(text_blocks),
    )


def get_slide_notes(document_path: Path, slide_number: int) -> str:
    presentation = _open_presentation(document_path)
    slide = _resolve_slide(presentation, slide_number)
    return _notes_text(slide)


def create_pptx(output_path: Path) -> Path:
    presentation = _open_empty_presentation()
    presentation.save(output_path)
    return output_path


def add_slide(document_path: Path, output_path: Path | None = None) -> tuple[Path, str]:
    presentation = _open_presentation(document_path)
    layout = _default_slide_layout(presentation)
    presentation.slides.add_slide(layout)
    slide_number = len(presentation.slides)
    target_path = _target_path(document_path, output_path)
    presentation.save(target_path)
    return target_path, f"pptx:slide:{slide_number}"


def add_textbox(
    document_path: Path,
    slide_locator: str,
    text: str,
    left: int | None = None,
    top: int | None = None,
    width: int | None = None,
    height: int | None = None,
    output_path: Path | None = None,
) -> tuple[Path, str]:
    presentation = _open_presentation(document_path)
    slide_number = _slide_number_from_any_locator(slide_locator)
    slide = _resolve_slide(presentation, slide_number)
    resolved_left, resolved_top, resolved_width, resolved_height = _default_textbox_geometry(
        presentation,
        left=left,
        top=top,
        width=width,
        height=height,
    )
    shape = slide.shapes.add_textbox(resolved_left, resolved_top, resolved_width, resolved_height)
    shape.text_frame.text = text
    locator = f"pptx:slide:{slide_number}:shape:{shape.shape_id}"
    target_path = _target_path(document_path, output_path)
    presentation.save(target_path)
    return target_path, locator


def style_run(
    document_path: Path,
    locator: str,
    style: InlineStyle,
    clear_fields: list[str] | tuple[str, ...],
    output_path: Path | None = None,
) -> tuple[Path, str, dict[str, object]]:
    presentation = _open_presentation(document_path)
    target = _resolve_text_target(presentation, locator, require_run=True)
    clear_set = _normalize_clear_fields(clear_fields, _INLINE_STYLE_FIELDS)
    skipped_fields = _apply_pptx_inline_style(target["run"], style, clear_set)
    target_path = _target_path(document_path, output_path)
    presentation.save(target_path)
    return target_path, target["shape_locator"], {"cleared_fields": clear_set, "skipped_fields": skipped_fields}


def style_paragraph(
    document_path: Path,
    locator: str,
    style: BlockStyle,
    clear_fields: list[str] | tuple[str, ...],
    output_path: Path | None = None,
) -> tuple[Path, str, dict[str, object]]:
    presentation = _open_presentation(document_path)
    target = _resolve_text_target(presentation, locator, require_run=False)
    clear_set = _normalize_clear_fields(clear_fields, _BLOCK_STYLE_FIELDS)
    skipped_fields = _apply_pptx_block_style(target["paragraph"], style, clear_set)
    target_path = _target_path(document_path, output_path)
    presentation.save(target_path)
    return target_path, target["shape_locator"], {"cleared_fields": clear_set, "skipped_fields": skipped_fields}


def parse_item_id(item_id: str) -> tuple[int, int]:
    parts = item_id.split(":")
    if len(parts) != 4 or parts[0] != "slide" or parts[2] != "shape":
        raise InvalidArgumentsError(f"Unsupported PPTX item id: {item_id}")

    try:
        slide_number = int(parts[1])
        shape_id = int(parts[3])
    except ValueError as exc:
        raise InvalidArgumentsError(f"Invalid PPTX item id: {item_id}") from exc

    if slide_number < 1:
        raise InvalidArgumentsError(f"Invalid PPTX slide number: {slide_number}")

    return slide_number, shape_id


def make_item_id(slide_number: int, shape_id: int) -> str:
    return f"slide:{slide_number}:shape:{shape_id}"


def _open_presentation(document_path: Path):
    if Presentation is None:
        raise RuntimeError("python-pptx is required for PPTX operations.")
    return Presentation(str(document_path))


def _open_empty_presentation():
    if Presentation is None:
        raise RuntimeError("python-pptx is required for PPTX operations.")
    return Presentation()


def _document_ref(document_path: Path) -> DocumentRef:
    resolved_path = document_path.resolve()
    stat = resolved_path.stat()
    return DocumentRef(
        document_id=resolved_path.as_posix(),
        path=resolved_path,
        file_type="pptx",
        display_name=resolved_path.name,
        modified_time=stat.st_mtime,
    )


def _resolve_shape(presentation, item_id: str):
    slide_number, shape_id = parse_item_id(item_id)
    slide = _resolve_slide(presentation, slide_number)

    for shape in slide.shapes:
        if shape.shape_id == shape_id:
            return shape

    raise TargetNotFoundError(f"Shape {shape_id} does not exist on slide {slide_number}.")


def _resolve_slide(presentation, slide_number: int):
    if slide_number < 1:
        raise InvalidArgumentsError(f"Invalid PPTX slide number: {slide_number}")

    try:
        return presentation.slides[slide_number - 1]
    except IndexError as exc:
        raise TargetNotFoundError(
            f"Slide {slide_number} does not exist in the presentation."
        ) from exc


def _require_text_frame(shape):
    if not getattr(shape, "has_text_frame", False):
        raise TargetNotEditableError("target not editable")
    return shape.text_frame


def _text_frame_text(text_frame) -> str:
    return "\n".join(paragraph.text for paragraph in text_frame.paragraphs)


def _shape_index(shape) -> int:
    return shape.element.getparent().index(shape.element)


def _target_path(document_path: Path, output_path: Path | None) -> Path:
    return document_path if output_path is None else output_path


def _default_slide_layout(presentation):
    if not presentation.slide_layouts:
        raise RuntimeError("Presentation has no slide layouts.")
    for layout in presentation.slide_layouts:
        if getattr(layout, "name", "").lower() == "blank":
            return layout
    if len(presentation.slide_layouts) > 6:
        return presentation.slide_layouts[6]
    return presentation.slide_layouts[-1]


def _slide_number_from_any_locator(locator: str) -> int:
    canonical = to_v2_locator(locator, file_type="pptx")
    parts = parse_locator(canonical).components
    if len(parts) >= 3 and parts[:2] == ("pptx", "slide"):
        return _parse_index(parts[2], locator, label="slide")
    raise InvalidArgumentsError(f"Unsupported PPTX slide locator: {locator}")


def _default_textbox_geometry(
    presentation,
    *,
    left: int | None,
    top: int | None,
    width: int | None,
    height: int | None,
) -> tuple[int, int, int, int]:
    slide_width = int(presentation.slide_width)
    slide_height = int(presentation.slide_height)
    resolved_width = width if width is not None else int(slide_width * 0.55)
    resolved_height = height if height is not None else int(slide_height * 0.2)
    resolved_left = left if left is not None else int((slide_width - resolved_width) / 2)
    resolved_top = top if top is not None else int((slide_height - resolved_height) / 2)
    return resolved_left, resolved_top, resolved_width, resolved_height


def _resolve_text_target(presentation, locator: str, *, require_run: bool) -> dict[str, object]:
    canonical = to_v2_locator(locator, file_type="pptx")
    parts = parse_locator(canonical).components
    if len(parts) < 5 or parts[:2] != ("pptx", "slide") or parts[3] != "shape":
        raise InvalidArgumentsError(f"Unsupported PPTX text locator: {locator}")

    slide_number = _parse_index(parts[2], locator, label="slide")
    shape_id = _parse_index(parts[4], locator, label="shape")
    shape = _resolve_shape(presentation, make_item_id(slide_number, shape_id))
    text_frame = _require_text_frame(shape)
    paragraph_index = 0
    run_index = 0
    if len(parts) >= 7:
        if parts[5] != "para":
            raise InvalidArgumentsError(f"Unsupported PPTX text locator: {locator}")
        paragraph_index = _parse_index(parts[6], locator, label="paragraph")
    try:
        paragraph = text_frame.paragraphs[paragraph_index]
    except IndexError as exc:
        raise TargetNotFoundError(
            f"Paragraph {paragraph_index} does not exist in PPTX shape {shape_id} on slide {slide_number}."
        ) from exc

    run = None
    if len(parts) >= 9:
        if parts[7] != "run":
            raise InvalidArgumentsError(f"Unsupported PPTX text locator: {locator}")
        run_index = _parse_index(parts[8], locator, label="run")
    if require_run:
        if not paragraph.runs:
            run = paragraph.add_run()
        else:
            try:
                run = paragraph.runs[run_index]
            except IndexError as exc:
                raise TargetNotFoundError(
                    f"Run {run_index} does not exist in PPTX paragraph {paragraph_index} on slide {slide_number}."
                ) from exc

    return {
        "canonical_locator": canonical,
        "shape_locator": f"pptx:slide:{slide_number}:shape:{shape_id}",
        "slide_number": slide_number,
        "shape_id": shape_id,
        "paragraph_index": paragraph_index,
        "run_index": run_index,
        "paragraph": paragraph,
        "run": run,
    }


def _slide_text_blocks(slide) -> list[SlideTextBlock]:
    blocks: list[SlideTextBlock] = []
    for position, shape in enumerate(slide.shapes):
        if not getattr(shape, "has_text_frame", False):
            continue
        text = _text_frame_text(shape.text_frame)
        blocks.append(
            SlideTextBlock(
                position=position,
                shape_id=shape.shape_id,
                shape_name=getattr(shape, "name", None),
                preview=text[:120],
                text=text,
                metadata={
                    "shape_index": position,
                    "is_placeholder": bool(getattr(shape, "is_placeholder", False)),
                },
            )
        )
    return blocks


def _slide_number_from_locator(locator: str) -> int:
    normalized = locator.strip()
    if normalized.startswith("slide:") and ":shape:" in normalized:
        slide_number, _ = parse_item_id(normalized)
        return slide_number
    parts = normalized.split(":")
    if len(parts) == 2 and parts[0] == "slide":
        try:
            slide_number = int(parts[1])
        except ValueError as exc:
            raise InvalidArgumentsError(f"Invalid slide locator: {locator}") from exc
        if slide_number < 1:
            raise InvalidArgumentsError(f"Invalid PPTX slide number: {slide_number}")
        return slide_number
    raise InvalidArgumentsError(f"Unsupported PPTX locator: {locator}")


def _first_text_shape_locator(document_path: Path, slide_number: int) -> str | None:
    bundle = get_slide_bundle(document_path, slide_number)
    if not bundle.text_blocks:
        return None
    return make_item_id(slide_number, bundle.text_blocks[0].shape_id)


def _notes_text(slide) -> str:
    notes_slide = getattr(slide, "notes_slide", None)
    if notes_slide is None:
        return ""
    text_frame = getattr(notes_slide, "notes_text_frame", None)
    if text_frame is None:
        return ""
    lines = [paragraph.text for paragraph in text_frame.paragraphs if paragraph.text.strip()]
    return "\n".join(lines)


_INLINE_STYLE_FIELDS = frozenset(
    {
        "bold",
        "italic",
        "underline",
        "strike",
        "font_name",
        "font_size",
        "font_color",
        "highlight",
    }
)
_BLOCK_STYLE_FIELDS = frozenset(
    {
        "alignment",
        "indent_level",
        "left_indent",
        "right_indent",
        "spacing_before",
        "spacing_after",
        "line_spacing",
        "wrap_text",
        "vertical_alignment",
        "fill_color",
        "number_format",
    }
)
_PPTX_ALIGNMENT_MAP = {
    "left": None if PP_ALIGN is None else PP_ALIGN.LEFT,
    "center": None if PP_ALIGN is None else PP_ALIGN.CENTER,
    "right": None if PP_ALIGN is None else PP_ALIGN.RIGHT,
    "justify": None if PP_ALIGN is None else PP_ALIGN.JUSTIFY,
}
_PPTX_VERTICAL_ALIGNMENT_MAP = {
    "top": None if MSO_ANCHOR is None else MSO_ANCHOR.TOP,
    "center": None if MSO_ANCHOR is None else MSO_ANCHOR.MIDDLE,
    "bottom": None if MSO_ANCHOR is None else MSO_ANCHOR.BOTTOM,
}


def _normalize_clear_fields(
    clear_fields: list[str] | tuple[str, ...],
    allowed: frozenset[str],
) -> tuple[str, ...]:
    normalized: list[str] = []
    seen: set[str] = set()
    for field_name in clear_fields:
        if field_name not in allowed:
            raise InvalidArgumentsError(f"Unknown style field in clear_fields: {field_name}")
        if field_name not in seen:
            normalized.append(field_name)
            seen.add(field_name)
    return tuple(normalized)


def _apply_pptx_inline_style(run, style: InlineStyle, clear_fields: tuple[str, ...]) -> list[str]:
    clear_set = set(clear_fields)
    font = run.font
    skipped_fields: list[str] = []

    if "bold" in clear_set:
        font.bold = None
    elif style.bold is not None:
        font.bold = style.bold

    if "italic" in clear_set:
        font.italic = None
    elif style.italic is not None:
        font.italic = style.italic

    if "underline" in clear_set:
        font.underline = None
    elif style.underline is not None:
        font.underline = style.underline

    if "strike" in clear_set:
        font.strike = None
    elif style.strike is not None:
        font.strike = style.strike

    if "font_name" in clear_set:
        font.name = None
    elif style.font_name is not None:
        font.name = style.font_name

    if "font_size" in clear_set:
        font.size = None
    elif style.font_size is not None:
        font.size = Pt(style.font_size)

    if "font_color" in clear_set:
        font.color.rgb = None
    elif style.font_color is not None:
        font.color.rgb = RGBColor.from_string(_normalize_hex_color(style.font_color))

    if style.highlight is not None or "highlight" in clear_set:
        skipped_fields.append("highlight")

    return skipped_fields


def _apply_pptx_block_style(paragraph, style: BlockStyle, clear_fields: tuple[str, ...]) -> list[str]:
    clear_set = set(clear_fields)
    skipped_fields: list[str] = []

    if "alignment" in clear_set:
        paragraph.alignment = None
    elif style.alignment is not None:
        paragraph.alignment = _pptx_alignment_value(style.alignment)

    if "indent_level" in clear_set:
        paragraph.level = 0
    elif style.indent_level is not None:
        paragraph.level = style.indent_level

    if "spacing_before" in clear_set:
        paragraph.space_before = None
    elif style.spacing_before is not None:
        paragraph.space_before = Pt(style.spacing_before)

    if "spacing_after" in clear_set:
        paragraph.space_after = None
    elif style.spacing_after is not None:
        paragraph.space_after = Pt(style.spacing_after)

    if "line_spacing" in clear_set:
        paragraph.line_spacing = None
    elif style.line_spacing is not None:
        paragraph.line_spacing = style.line_spacing

    for field_name in ("left_indent", "right_indent", "fill_color", "number_format"):
        if getattr(style, field_name) is not None or field_name in clear_set:
            skipped_fields.append(field_name)

    if style.wrap_text is not None or "wrap_text" in clear_set:
        skipped_fields.append("wrap_text")

    if style.vertical_alignment is not None or "vertical_alignment" in clear_set:
        skipped_fields.append("vertical_alignment")

    return skipped_fields


def _pptx_alignment_value(raw: str):
    normalized = raw.strip().lower()
    if normalized not in _PPTX_ALIGNMENT_MAP:
        raise InvalidArgumentsError(f"Unsupported PPTX alignment: {raw}")
    return _PPTX_ALIGNMENT_MAP[normalized]


def _normalize_hex_color(value: str) -> str:
    normalized = value.strip().lstrip("#").upper()
    if len(normalized) != 6 or any(character not in "0123456789ABCDEF" for character in normalized):
        raise InvalidArgumentsError(f"Invalid RGB hex color: {value}")
    return normalized


def _parse_index(raw: str, locator: str, *, label: str) -> int:
    try:
        return int(raw)
    except ValueError as exc:
        raise InvalidArgumentsError(f"Invalid PPTX {label} index in locator: {locator}") from exc
