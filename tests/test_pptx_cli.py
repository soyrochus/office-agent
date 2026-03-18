from __future__ import annotations

import subprocess
import sys
from pathlib import Path

from pptx import Presentation

from offagent.adapters import pptx_adapter


def _run_cli(config_path, *args: str) -> subprocess.CompletedProcess[str]:
    return subprocess.run(
        [
            sys.executable,
            "-m",
            "offagent",
            *args,
            "--config",
            str(config_path),
        ],
        capture_output=True,
        text=True,
        check=False,
    )


def _parse_output_path(stdout: str) -> Path:
    return Path(stdout.strip().split("\t")[-1])


def test_pptx_cli_round_trip(sample_pptx, config_path) -> None:
    extracted_items = pptx_adapter.extract_document(sample_pptx)
    summary_item = next(item for item in extracted_items if "Supplier shall present" in item.content_text)
    editable_item = next(item for item in extracted_items if item.metadata["slide_number"] == 2)

    index_result = _run_cli(config_path, "index", str(sample_pptx))
    search_result = _run_cli(config_path, "search", "Supplier shall", "--type", "pptx")
    locate_slide_result = _run_cli(config_path, "locate", "--doc", str(sample_pptx), "--slide", "1")
    locate_shape_result = _run_cli(
        config_path,
        "locate",
        "--doc",
        str(sample_pptx),
        "--slide",
        "1",
        "--shape",
        str(summary_item.metadata["shape_id"]),
    )
    read_result = _run_cli(config_path, "read", "--doc", str(sample_pptx), "--item", summary_item.item_id)
    replace_result = _run_cli(
        config_path,
        "replace",
        "--doc",
        str(sample_pptx),
        "--item",
        editable_item.item_id,
        "--text",
        "CLI updated speaker notes.",
    )
    replace_output_path = _parse_output_path(replace_result.stdout)
    append_result = _run_cli(
        config_path,
        "append",
        "--doc",
        str(replace_output_path),
        "--item",
        editable_item.item_id,
        "--text",
        "\nCLI appended action.",
    )
    append_output_path = _parse_output_path(append_result.stdout)

    assert index_result.returncode == 0, index_result.stderr
    assert "indexed 1" in index_result.stdout
    assert search_result.returncode == 0, search_result.stderr
    assert summary_item.item_id in search_result.stdout
    assert locate_slide_result.returncode == 0, locate_slide_result.stderr
    assert (
        len(
            [
                line
                for line in locate_slide_result.stdout.splitlines()
                if line.startswith("slide:1:shape:")
            ]
        )
        == 3
    )
    assert locate_shape_result.returncode == 0, locate_shape_result.stderr
    assert summary_item.item_id in locate_shape_result.stdout
    assert read_result.returncode == 0, read_result.stderr
    assert "Supplier shall present the rollout plan." in read_result.stdout
    assert replace_result.returncode == 0, replace_result.stderr
    assert append_result.returncode == 0, append_result.stderr

    original_presentation = Presentation(str(sample_pptx))
    presentation = Presentation(str(append_output_path))
    original_shape = next(shape for shape in original_presentation.slides[1].shapes if shape.has_text_frame)
    editable_shape = next(shape for shape in presentation.slides[1].shapes if shape.has_text_frame)
    assert original_shape.text_frame.text == "Editable speaker notes"
    assert editable_shape.text_frame.text == "CLI updated speaker notes.\nCLI appended action."
