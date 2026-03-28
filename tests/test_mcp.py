from __future__ import annotations

import asyncio
import sys
from pathlib import Path

from docx import Document
from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client
from openpyxl import load_workbook
from pptx import Presentation


def _server_params(config_path: Path) -> StdioServerParameters:
    return StdioServerParameters(
        command=sys.executable,
        args=["-m", "offagent", "mcp", "--config", str(config_path)],
    )


def _run_mcp(config_path: Path, callback):
    async def runner():
        async with stdio_client(_server_params(config_path)) as (read_stream, write_stream):
            async with ClientSession(read_stream, write_stream) as session:
                await session.initialize()
                return await callback(session)

    return asyncio.run(runner())


async def _call_tool(session: ClientSession, name: str, arguments: dict | None = None):
    result = await session.call_tool(name, arguments=arguments)
    assert not result.isError, _tool_error_text(result)
    return result.structuredContent


def _tool_error_text(result) -> str:
    parts: list[str] = []
    for content in result.content:
        text = getattr(content, "text", None)
        if text is not None:
            parts.append(text)
    return "\n".join(parts)


def test_mcp_lists_tools_and_runs_docx_flow(sample_docx, config_path) -> None:
    def scenario(session: ClientSession):
        async def run():
            tools = await session.list_tools()
            tool_map = {tool.name: tool for tool in tools.tools}
            assert set(tool_map) == {
                "index_documents",
                "refresh_document",
                "list_documents",
                "search_documents",
                "get_structure",
                "get_section",
                "get_node",
                "write_node",
                "insert_content",
                "xlsx_insert_rows",
                "docx_get_tables",
            }
            assert all(tool.inputSchema is not None for tool in tool_map.values())
            assert all(tool.outputSchema is not None for tool in tool_map.values())

            index_result = await _call_tool(
                session,
                "index_documents",
                {"paths": [str(sample_docx)]},
            )
            assert index_result["files_indexed"] == 1

            documents_result = await _call_tool(session, "list_documents")
            document = documents_result["documents"][0]
            assert document["path"] == str(sample_docx)

            search_result = await _call_tool(
                session,
                "search_documents",
                {"query": "Supplier shall", "file_type": "docx"},
            )
            hit = search_result["hits"][0]
            assert hit["item_id"] == "para:3"

            node_result = await _call_tool(
                session,
                "get_node",
                {"document_id": document["document_id"], "node_id": hit["locator"]},
            )
            assert node_result["text"] == "Supplier shall deliver by Friday."

            structure = await _call_tool(
                session,
                "get_structure",
                {"document_id": document["document_id"]},
            )
            assert [section["section_type"] for section in structure["sections"]] == [
                "paragraph",
                "paragraph",
                "paragraph",
                "paragraph",
                "table",
            ]

            table_locator = structure["sections"][-1]["locator"]
            table_section = await _call_tool(
                session,
                "get_section",
                {"document_id": document["document_id"], "section_id": table_locator},
            )
            assert table_section["rows"] == [["Table text to ignore"]]

            replace_result = await _call_tool(
                session,
                "write_node",
                {
                    "document_id": document["document_id"],
                    "node_id": "para:1",
                    "content": "MCP replaced text.",
                },
            )
            blank_node = await _call_tool(
                session,
                "get_node",
                {"document_id": replace_result["document_id"], "node_id": "para:2"},
            )
            append_result = await _call_tool(
                session,
                "write_node",
                {
                    "document_id": replace_result["document_id"],
                    "node_id": "para:2",
                    "content": f'{blank_node["text"]}MCP appended text.',
                },
            )
            insert_result = await _call_tool(
                session,
                "insert_content",
                {
                    "document_id": append_result["document_id"],
                    "content": "Inserted after refactor.",
                    "style_name": "Heading 2",
                },
            )

            refresh_result = await _call_tool(
                session,
                "refresh_document",
                {"document_id": document["document_id"]},
            )
            assert refresh_result["files_indexed"] == 1

            return insert_result

        return run()

    insert_result = _run_mcp(config_path, scenario)

    original_document = Document(str(sample_docx))
    final_document = Document(insert_result["output_path"])
    assert original_document.paragraphs[1].text == "Alpha paragraph for search."
    assert final_document.paragraphs[1].text == "MCP replaced text."
    assert final_document.paragraphs[1].runs[0].bold is True
    assert final_document.paragraphs[2].text == "MCP appended text."
    assert final_document.paragraphs[-1].text == "Inserted after refactor."
    assert final_document.paragraphs[-1].style.name == "Heading 2"


def test_mcp_pptx_and_xlsx_round_trips(sample_pptx, sample_xlsx, config_path) -> None:
    def scenario(session: ClientSession):
        async def run():
            await _call_tool(session, "index_documents", {"paths": [str(sample_pptx), str(sample_xlsx)]})

            documents_result = await _call_tool(session, "list_documents")
            docs_by_type = {document["file_type"]: document for document in documents_result["documents"]}

            pptx_search = await _call_tool(
                session,
                "search_documents",
                {"query": "Supplier shall", "file_type": "pptx"},
            )
            pptx_hit = pptx_search["hits"][0]
            pptx_read = await _call_tool(
                session,
                "get_node",
                {"document_id": docs_by_type["pptx"]["document_id"], "node_id": pptx_hit["item_id"]},
            )
            assert "Supplier shall present the rollout plan." in pptx_read["text"]

            editable_search = await _call_tool(
                session,
                "search_documents",
                {"query": "Editable speaker notes", "file_type": "pptx"},
            )
            editable_item_id = editable_search["hits"][0]["item_id"]
            pptx_replace = await _call_tool(
                session,
                "write_node",
                {
                    "document_id": docs_by_type["pptx"]["document_id"],
                    "node_id": editable_item_id,
                    "content": "MCP updated speaker notes.",
                },
            )
            pptx_append = await _call_tool(
                session,
                "write_node",
                {
                    "document_id": pptx_replace["document_id"],
                    "node_id": editable_item_id,
                    "content": "MCP updated speaker notes.\nMCP appended action.",
                },
            )

            xlsx_structure = await _call_tool(
                session,
                "get_structure",
                {"document_id": docs_by_type["xlsx"]["document_id"]},
            )
            notes_locator = next(
                section["locator"]
                for section in xlsx_structure["sections"]
                if section["sheet_name"] == "Notes 2026"
            )
            xlsx_section = await _call_tool(
                session,
                "get_section",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "section_id": notes_locator,
                    "cell_range": "A1:A2",
                },
            )
            assert [cell["coordinate"] for cell in xlsx_section["cells"]] == ["A1", "A2"]

            xlsx_write = await _call_tool(
                session,
                "write_node",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "node_id": "sheet:Budget2026!B2",
                    "content": "125001",
                },
            )
            note_node = await _call_tool(
                session,
                "get_node",
                {
                    "document_id": xlsx_write["document_id"],
                    "node_id": "sheet:Notes 2026!A2",
                },
            )
            xlsx_append = await _call_tool(
                session,
                "write_node",
                {
                    "document_id": xlsx_write["document_id"],
                    "node_id": "sheet:Notes 2026!A2",
                    "content": f'{note_node["text"]} MCP added note.',
                },
            )
            inserted_rows = await _call_tool(
                session,
                "xlsx_insert_rows",
                {
                    "document_id": xlsx_append["document_id"],
                    "sheet_name": "Notes 2026",
                    "rows": [["Row note", "Finance"]],
                },
            )

            return {
                "pptx_output": pptx_append["output_path"],
                "xlsx_output": inserted_rows["output_path"],
            }

        return run()

    outputs = _run_mcp(config_path, scenario)

    presentation = Presentation(outputs["pptx_output"])
    editable_shape = next(shape for shape in presentation.slides[1].shapes if shape.has_text_frame)
    assert "MCP updated speaker notes." in editable_shape.text_frame.text
    assert "MCP appended action." in editable_shape.text_frame.text

    workbook = load_workbook(outputs["xlsx_output"])
    assert workbook["Budget2026"]["B2"].value == 125001
    assert workbook["Notes 2026"]["A2"].value == "Follow up with finance. MCP added note."
    assert workbook["Notes 2026"]["A3"].value == "Row note"
    assert workbook["Notes 2026"]["B3"].value == "Finance"


def test_mcp_returns_structured_errors(sample_docx, sample_xlsx, config_path) -> None:
    def scenario(session: ClientSession):
        async def run():
            await _call_tool(session, "index_documents", {"paths": [str(sample_docx), str(sample_xlsx)]})
            documents_result = await _call_tool(session, "list_documents")
            docs_by_path = {document["path"]: document for document in documents_result["documents"]}

            unsupported = await session.call_tool(
                "insert_content",
                {"document_id": docs_by_path[str(sample_xlsx)]["document_id"], "content": "blocked"},
            )
            assert unsupported.isError
            assert "requires a .docx document" in _tool_error_text(unsupported)

            workbook = load_workbook(sample_xlsx)
            del workbook["Budget2026"]
            workbook.save(sample_xlsx)

            stale_result = await session.call_tool(
                "write_node",
                {
                    "document_id": docs_by_path[str(sample_xlsx)]["document_id"],
                    "node_id": "sheet:Budget2026!B2",
                    "content": "125001",
                },
            )
            assert stale_result.isError
            assert "stale locator" in _tool_error_text(stale_result)

        return run()

    _run_mcp(config_path, scenario)
