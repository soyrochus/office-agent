from __future__ import annotations

from docx import Document
from openpyxl import load_workbook
from pptx import Presentation
import pytest

from offagent.app.services import AppServices
from offagent.config import AppConfig
from offagent.errors import TargetNotEditableError


def _services(tmp_path) -> AppServices:
    return AppServices(
        AppConfig(index_path=tmp_path / "state" / "index.sqlite3", document_roots=(tmp_path,))
    )


def test_create_object_appends_docx_paragraph(sample_docx, tmp_path) -> None:
    services = _services(tmp_path)
    document = services.index_document(sample_docx)

    result = services.create_object(
        document.document_id,
        "docx:document",
        "paragraph",
        {"text": "New appendix paragraph."},
    )

    updated = Document(str(result.output_path))
    assert result.object_type == "paragraph"
    assert result.locator == "docx:para:4"
    assert updated.paragraphs[-1].text == "New appendix paragraph."


def test_update_object_uses_existing_write_paths(sample_docx, sample_pptx, sample_xlsx, tmp_path) -> None:
    services = _services(tmp_path)
    docx_document = services.index_document(sample_docx)
    pptx_document = services.index_document(sample_pptx)
    xlsx_document = services.index_document(sample_xlsx)

    docx_result = services.update_object(
        docx_document.document_id,
        "docx:para:1",
        {"text": "Updated alpha paragraph."},
    )
    pptx_shape = next(
        child.locator
        for child in services.list_children(pptx_document.document_id, "pptx:slide:2")
        if child.object_type == "text_shape"
    )
    pptx_result = services.update_object(
        pptx_document.document_id,
        pptx_shape,
        {"text": "Updated speaker notes."},
    )
    xlsx_result = services.update_object(
        xlsx_document.document_id,
        "xlsx:sheet:Budget2026!B2",
        {"value": "125001"},
    )

    updated_docx = Document(str(docx_result.output_path))
    updated_pptx = Presentation(str(pptx_result.output_path))
    updated_xlsx = load_workbook(xlsx_result.output_path)

    assert updated_docx.paragraphs[1].text == "Updated alpha paragraph."
    assert next(shape for shape in updated_pptx.slides[1].shapes if shape.has_text_frame).text == "Updated speaker notes."
    assert updated_xlsx["Budget2026"]["B2"].value == 125001


def test_move_object_reorders_pptx_slides(sample_pptx, tmp_path) -> None:
    services = _services(tmp_path)
    document = services.index_document(sample_pptx)

    result = services.move_object(
        document.document_id,
        "pptx:slide:2",
        "pptx:presentation",
        1,
    )

    presentation = Presentation(str(result.output_path))
    first_text_shape = next(shape for shape in presentation.slides[0].shapes if shape.has_text_frame)
    assert result.locator == "pptx:slide:1"
    assert first_text_shape.text == "Editable speaker notes"


def test_copy_object_duplicates_pptx_slide(sample_pptx, tmp_path) -> None:
    services = _services(tmp_path)
    document = services.index_document(sample_pptx)

    result = services.copy_object(
        document.document_id,
        "pptx:slide:1",
        "pptx:presentation",
        3,
    )

    presentation = Presentation(str(result.output_path))
    third_slide_texts = [shape.text for shape in presentation.slides[2].shapes if shape.has_text_frame]
    assert result.locator == "pptx:slide:3"
    assert len(presentation.slides) == 3
    assert any("Quarterly Planning" in text for text in third_slide_texts)


def test_batch_edit_supports_dry_run_and_atomic_failure(sample_docx, tmp_path) -> None:
    services = _services(tmp_path)
    document = services.index_document(sample_docx)

    dry_run = services.batch_edit(
        document.document_id,
        [
            {
                "op": "update_object",
                "locator": "docx:para:1",
                "properties": {"text": "Batch updated alpha."},
            },
            {
                "op": "create_object",
                "parent_locator": "docx:document",
                "object_type": "paragraph",
                "properties": {"text": "Batch appendix."},
            },
        ],
        dry_run=True,
    )

    assert dry_run.dry_run is True
    assert dry_run.output_path is None
    assert Document(str(sample_docx)).paragraphs[1].text == "Alpha paragraph for search."

    with pytest.raises(TargetNotEditableError):
        services.batch_edit(
            document.document_id,
            [
                {
                    "op": "update_object",
                    "locator": "docx:para:1",
                    "properties": {"text": "Should not persist."},
                },
                {
                    "op": "update_object",
                    "locator": "docx:document",
                    "properties": {"text": "Invalid target."},
                },
            ],
        )

    assert Document(str(sample_docx)).paragraphs[1].text == "Alpha paragraph for search."


def test_batch_edit_applies_multiple_operations(sample_docx, tmp_path) -> None:
    services = _services(tmp_path)
    document = services.index_document(sample_docx)

    result = services.batch_edit(
        document.document_id,
        [
            {
                "op": "update_object",
                "locator": "docx:para:1",
                "properties": {"text": "Batch updated alpha."},
            },
            {
                "op": "create_object",
                "parent_locator": "docx:document",
                "object_type": "paragraph",
                "properties": {"text": "Batch appendix."},
            },
        ],
    )

    updated = Document(str(result.output_path))
    assert result.dry_run is False
    assert len(result.operations) == 2
    assert updated.paragraphs[1].text == "Batch updated alpha."
    assert updated.paragraphs[-1].text == "Batch appendix."


def test_delete_object_removes_supported_targets_and_enforces_capability(sample_docx, sample_pptx, sample_xlsx, tmp_path) -> None:
    services = _services(tmp_path)
    docx_document = services.index_document(sample_docx)
    pptx_document = services.index_document(sample_pptx)
    xlsx_document = services.index_document(sample_xlsx)

    docx_result = services.delete_object(docx_document.document_id, "docx:para:2")
    pptx_result = services.delete_object(pptx_document.document_id, "pptx:slide:2")
    xlsx_result = services.delete_object(xlsx_document.document_id, "xlsx:sheet:Notes 2026")

    updated_docx = Document(str(docx_result.output_path))
    updated_pptx = Presentation(str(pptx_result.output_path))
    updated_xlsx = load_workbook(xlsx_result.output_path)

    assert len(updated_docx.paragraphs) == 3
    assert all(paragraph.text != "" for paragraph in updated_docx.paragraphs)
    assert len(updated_pptx.slides) == 1
    assert updated_xlsx.sheetnames == ["Budget2026"]

    with pytest.raises(TargetNotEditableError, match="does not support delete"):
        services.delete_object(docx_document.document_id, "docx:document")


def test_escape_hatch_service_methods_round_trip(sample_docx, sample_pptx, sample_xlsx, tmp_path) -> None:
    services = _services(tmp_path)
    docx_document = services.index_document(sample_docx)
    pptx_document = services.index_document(sample_pptx)
    xlsx_document = services.index_document(sample_xlsx)

    styled = services.docx_set_paragraph_style(docx_document.document_id, "docx:para:1", "Heading 2")
    page_break = services.docx_insert_page_break(styled.document_id, "docx:para:1")
    table = services.docx_add_table(page_break.document_id, 2, 2, style_name="Table Grid")
    merged = services.docx_merge_table_cells(
        table.document_id,
        f"{table.locator}:row:0:cell:0",
        f"{table.locator}:row:0:cell:1",
    )

    added_slide = services.pptx_add_slide(pptx_document.document_id, layout_index=6)
    duplicated_slide = services.pptx_duplicate_slide(pptx_document.document_id, "pptx:slide:1", position=2)
    relayout = services.pptx_set_slide_layout(pptx_document.document_id, "pptx:slide:2", layout_index=5)
    text_shape = services.pptx_add_text_shape(
        pptx_document.document_id,
        "pptx:slide:2",
        "Service textbox",
        left=914400,
        top=914400,
        width=1828800,
        height=914400,
    )

    written_range = services.xlsx_write_range(
        xlsx_document.document_id,
        "xlsx:sheet:Budget2026!A1:B2",
        [["Header", "Budget"], ["Q1", "200"]],
    )
    inserted_rows = services.xlsx_insert_rows_at(xlsx_document.document_id, "xlsx:sheet:Budget2026", 2, 1)
    inserted_columns = services.xlsx_insert_columns(xlsx_document.document_id, "xlsx:sheet:Budget2026", 1, 1)
    formula = services.xlsx_set_formula(xlsx_document.document_id, "xlsx:sheet:Budget2026!D4", "=SUM(1,2)")
    merged_cells = services.xlsx_merge_cells(xlsx_document.document_id, "xlsx:sheet:Budget2026!A4:B4")

    assert Document(str(styled.output_path)).paragraphs[1].style.name == "Heading 2"
    assert page_break.object_type == "page_break"
    assert merged.locator == f"{table.locator}:row:0:cell:0"

    assert len(Presentation(str(added_slide.output_path)).slides) == 3
    assert duplicated_slide.locator == "pptx:slide:2"
    assert Presentation(str(relayout.output_path)).slides[1].slide_layout.name == Presentation(str(relayout.output_path)).slide_layouts[5].name
    assert text_shape.object_type == "text_shape"

    assert load_workbook(written_range.output_path)["Budget2026"]["B2"].value == 200
    assert load_workbook(inserted_rows.output_path)["Budget2026"]["B3"].value == 125000
    assert load_workbook(inserted_columns.output_path)["Budget2026"]["B1"].value == "Quarterly Budget"
    assert formula.locator == "xlsx:sheet:Budget2026:formula_cell:D4"
    assert "A4:B4" in {str(cell_range) for cell_range in load_workbook(merged_cells.output_path)["Budget2026"].merged_cells.ranges}
