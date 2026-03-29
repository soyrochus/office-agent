from __future__ import annotations

import asyncio
import sys
from pathlib import Path

from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client
from openpyxl import load_workbook
from pptx import Presentation
from docx import Document


def _server_params(config_path: Path) -> StdioServerParameters:
    return StdioServerParameters(
        command=sys.executable,
        args=["-m", "offagent", "mcp", "--config", str(config_path)],
    )


def _run_mcp(config_path: Path, callback):
    async def runner():
        async with stdio_client(_server_params(config_path)) as (read_stream, write_stream):
            async with ClientSession(read_stream, write_stream) as session:
                await session.initialize()
                return await callback(session)

    return asyncio.run(runner())


async def _call_tool(session: ClientSession, name: str, arguments: dict | None = None):
    result = await session.call_tool(name, arguments=arguments)
    assert not result.isError
    return result.structuredContent


def test_mcp_v2_tools_register_and_round_trip_docx(sample_docx, config_path) -> None:
    def scenario(session: ClientSession):
        async def run():
            tools = await session.list_tools()
            tool_map = {tool.name: tool for tool in tools.tools}
            assert {"search_objects", "get_object", "list_children"} <= set(tool_map)

            await _call_tool(session, "index_documents", {"paths": [str(sample_docx)]})
            documents = await _call_tool(session, "list_documents")
            document = documents["documents"][0]

            search = await _call_tool(
                session,
                "search_objects",
                {"query": "Supplier shall", "file_type": "docx"},
            )
            hit = search["hits"][0]
            assert hit["object_type"] == "paragraph"
            assert hit["match_mode"] == "keyword"

            obj = await _call_tool(
                session,
                "get_object",
                {"document_id": document["document_id"], "locator": hit["locator"]},
            )
            children = await _call_tool(
                session,
                "list_children",
                {
                    "document_id": document["document_id"],
                    "locator": "docx:document",
                    "child_type": "paragraph",
                    "limit": 2,
                },
            )
            return {"object": obj, "children": children}

        return run()

    result = _run_mcp(config_path, scenario)

    assert result["object"]["object_type"] == "paragraph"
    assert result["object"]["properties"]["text"] == "Supplier shall deliver by Friday."
    assert len(result["children"]["children"]) == 2
    assert result["children"]["children"][0]["object_type"] == "paragraph"


def test_mcp_v2_get_object_and_list_children_across_formats(
    sample_docx,
    sample_pptx,
    sample_xlsx,
    config_path,
) -> None:
    def scenario(session: ClientSession):
        async def run():
            await _call_tool(session, "index_documents", {"paths": [str(sample_docx), str(sample_pptx), str(sample_xlsx)]})
            documents = await _call_tool(session, "list_documents")
            docs_by_type = {document["file_type"]: document for document in documents["documents"]}

            docx_object = await _call_tool(
                session,
                "get_object",
                {"document_id": docs_by_type["docx"]["document_id"], "locator": "docx:para:1"},
            )
            pptx_children = await _call_tool(
                session,
                "list_children",
                {"document_id": docs_by_type["pptx"]["document_id"], "locator": "pptx:slide:1"},
            )
            xlsx_children = await _call_tool(
                session,
                "list_children",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "locator": "xlsx:sheet:Notes 2026",
                    "child_type": "row",
                },
            )
            return {"docx": docx_object, "pptx": pptx_children, "xlsx": xlsx_children}

        return run()

    result = _run_mcp(config_path, scenario)

    assert result["docx"]["object_type"] == "paragraph"
    assert any(child["object_type"] == "text_shape" for child in result["pptx"]["children"])
    assert [child["locator"] for child in result["xlsx"]["children"]] == [
        "xlsx:sheet:Notes 2026:row:1",
        "xlsx:sheet:Notes 2026:row:2",
    ]


def test_mcp_v2_mutation_tools_round_trip(sample_docx, sample_pptx, sample_xlsx, config_path) -> None:
    def scenario(session: ClientSession):
        async def run():
            await _call_tool(session, "index_documents", {"paths": [str(sample_docx), str(sample_pptx), str(sample_xlsx)]})
            documents = await _call_tool(session, "list_documents")
            docs_by_type = {document["file_type"]: document for document in documents["documents"]}

            created = await _call_tool(
                session,
                "create_object",
                {
                    "document_id": docs_by_type["docx"]["document_id"],
                    "parent_locator": "docx:document",
                    "object_type": "paragraph",
                    "properties": {"text": "Created from MCP V2."},
                },
            )
            updated = await _call_tool(
                session,
                "update_object",
                {
                    "document_id": created["document_id"],
                    "locator": "docx:para:1",
                    "properties": {"text": "Updated from MCP V2."},
                },
            )
            pptx_shape = next(
                child["locator"]
                for child in (
                    await _call_tool(
                        session,
                        "list_children",
                        {"document_id": docs_by_type["pptx"]["document_id"], "locator": "pptx:slide:2"},
                    )
                )["children"]
                if child["object_type"] == "text_shape"
            )
            pptx_updated = await _call_tool(
                session,
                "update_object",
                {
                    "document_id": docs_by_type["pptx"]["document_id"],
                    "locator": pptx_shape,
                    "properties": {"text": "Updated from MCP V2 PPTX."},
                },
            )
            deleted = await _call_tool(
                session,
                "delete_object",
                {
                    "document_id": updated["document_id"],
                    "locator": "docx:para:2",
                },
            )
            copied = await _call_tool(
                session,
                "copy_object",
                {
                    "document_id": docs_by_type["pptx"]["document_id"],
                    "locator": "pptx:slide:1",
                    "target_parent_locator": "pptx:presentation",
                    "position": 3,
                },
            )
            moved = await _call_tool(
                session,
                "move_object",
                {
                    "document_id": copied["document_id"],
                    "locator": "pptx:slide:2",
                    "new_parent_locator": "pptx:presentation",
                    "position": 1,
                },
            )
            xlsx_updated = await _call_tool(
                session,
                "update_object",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "locator": "xlsx:sheet:Budget2026!B2",
                    "properties": {"value": "125001"},
                },
            )
            xlsx_deleted = await _call_tool(
                session,
                "delete_object",
                {
                    "document_id": xlsx_updated["document_id"],
                    "locator": "xlsx:sheet:Notes 2026",
                },
            )
            batched = await _call_tool(
                session,
                "batch_edit",
                {
                    "document_id": deleted["document_id"],
                    "operations": [
                        {
                            "op": "update_object",
                            "locator": "docx:para:1",
                            "properties": {"text": "Batch updated from MCP V2."},
                        },
                        {
                            "op": "create_object",
                            "parent_locator": "docx:document",
                            "object_type": "paragraph",
                            "properties": {"text": "Batch created from MCP V2."},
                        },
                    ],
                },
            )
            return {
                "created": created,
                "updated": updated,
                "pptx_updated": pptx_updated,
                "deleted": deleted,
                "copied": copied,
                "moved": moved,
                "xlsx_updated": xlsx_updated,
                "xlsx_deleted": xlsx_deleted,
                "batched": batched,
            }

        return run()

    result = _run_mcp(config_path, scenario)

    assert result["created"]["locator"] == "docx:para:4"
    assert result["updated"]["locator"] == "docx:para:1"
    assert result["pptx_updated"]["object_type"] == "text_shape"
    assert result["deleted"]["locator"] is None
    assert result["copied"]["locator"] == "pptx:slide:3"
    assert result["moved"]["locator"] == "pptx:slide:1"
    assert result["xlsx_updated"]["locator"] == "xlsx:sheet:Budget2026!B2"
    assert result["xlsx_deleted"]["object_type"] == "worksheet"
    assert len(result["batched"]["operations"]) == 2


def test_mcp_v2_batch_edit_supports_dry_run_and_atomic_failure(sample_docx, config_path) -> None:
    def scenario(session: ClientSession):
        async def run():
            await _call_tool(session, "index_documents", {"paths": [str(sample_docx)]})
            document = (await _call_tool(session, "list_documents"))["documents"][0]

            dry_run = await _call_tool(
                session,
                "batch_edit",
                {
                    "document_id": document["document_id"],
                    "dry_run": True,
                    "operations": [
                        {
                            "op": "update_object",
                            "locator": "docx:para:1",
                            "properties": {"text": "Dry run update."},
                        }
                    ],
                },
            )
            failure = await session.call_tool(
                "batch_edit",
                {
                    "document_id": document["document_id"],
                    "operations": [
                        {
                            "op": "update_object",
                            "locator": "docx:para:1",
                            "properties": {"text": "Should not persist."},
                        },
                        {
                            "op": "update_object",
                            "locator": "docx:document",
                            "properties": {"text": "Invalid target."},
                        },
                    ],
                },
            )
            return {"dry_run": dry_run, "failure": failure}

        return run()

    result = _run_mcp(config_path, scenario)

    assert result["dry_run"]["dry_run"] is True
    assert result["dry_run"]["output_path"] is None
    assert result["failure"].isError
    assert Document(str(sample_docx)).paragraphs[1].text == "Alpha paragraph for search."


def test_mcp_v2_escape_hatch_tools_round_trip(sample_docx, sample_pptx, sample_xlsx, config_path) -> None:
    def scenario(session: ClientSession):
        async def run():
            await _call_tool(session, "index_documents", {"paths": [str(sample_docx), str(sample_pptx), str(sample_xlsx)]})
            documents = await _call_tool(session, "list_documents")
            docs_by_type = {document["file_type"]: document for document in documents["documents"]}

            docx_style = await _call_tool(
                session,
                "docx_set_paragraph_style",
                {
                    "document_id": docs_by_type["docx"]["document_id"],
                    "locator": "docx:para:1",
                    "style_name": "Heading 2",
                },
            )
            docx_break = await _call_tool(
                session,
                "docx_insert_page_break",
                {
                    "document_id": docs_by_type["docx"]["document_id"],
                    "locator": "docx:para:1",
                },
            )
            docx_table = await _call_tool(
                session,
                "docx_add_table",
                {
                    "document_id": docs_by_type["docx"]["document_id"],
                    "row_count": 2,
                    "column_count": 2,
                    "style_name": "Table Grid",
                },
            )
            docx_merged = await _call_tool(
                session,
                "docx_merge_table_cells",
                {
                    "document_id": docx_table["document_id"],
                    "start_locator": f'{docx_table["locator"]}:row:0:cell:0',
                    "end_locator": f'{docx_table["locator"]}:row:0:cell:1',
                },
            )

            pptx_added = await _call_tool(
                session,
                "pptx_add_slide",
                {"document_id": docs_by_type["pptx"]["document_id"], "layout_index": 6},
            )
            pptx_duplicated = await _call_tool(
                session,
                "pptx_duplicate_slide",
                {
                    "document_id": docs_by_type["pptx"]["document_id"],
                    "locator": "pptx:slide:1",
                    "position": 2,
                },
            )
            pptx_layout = await _call_tool(
                session,
                "pptx_set_slide_layout",
                {
                    "document_id": docs_by_type["pptx"]["document_id"],
                    "locator": "pptx:slide:2",
                    "layout_index": 5,
                },
            )
            pptx_text = await _call_tool(
                session,
                "pptx_add_text_shape",
                {
                    "document_id": docs_by_type["pptx"]["document_id"],
                    "locator": "pptx:slide:2",
                    "text": "Inserted textbox",
                    "left": 914400,
                    "top": 914400,
                    "width": 1828800,
                    "height": 914400,
                },
            )

            xlsx_write = await _call_tool(
                session,
                "xlsx_write_range",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "locator": "xlsx:sheet:Budget2026!A1:B2",
                    "values": [["Header", "Budget"], ["Q1", "200"]],
                },
            )
            xlsx_insert_rows = await _call_tool(
                session,
                "xlsx_insert_rows",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "locator": "xlsx:sheet:Budget2026",
                    "row_number": 2,
                    "count": 1,
                },
            )
            xlsx_insert_columns = await _call_tool(
                session,
                "xlsx_insert_columns",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "locator": "xlsx:sheet:Budget2026",
                    "column_index": 1,
                    "count": 1,
                },
            )
            xlsx_formula = await _call_tool(
                session,
                "xlsx_set_formula",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "locator": "xlsx:sheet:Budget2026!D4",
                    "formula": "=SUM(1,2)",
                },
            )
            xlsx_merged = await _call_tool(
                session,
                "xlsx_merge_cells",
                {
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "locator": "xlsx:sheet:Budget2026!A4:B4",
                },
            )

            return {
                "docx_style": docx_style,
                "docx_break": docx_break,
                "docx_table": docx_table,
                "docx_merged": docx_merged,
                "pptx_added": pptx_added,
                "pptx_duplicated": pptx_duplicated,
                "pptx_layout": pptx_layout,
                "pptx_text": pptx_text,
                "xlsx_write": xlsx_write,
                "xlsx_insert_rows": xlsx_insert_rows,
                "xlsx_insert_columns": xlsx_insert_columns,
                "xlsx_formula": xlsx_formula,
                "xlsx_merged": xlsx_merged,
            }

        return run()

    result = _run_mcp(config_path, scenario)

    assert Document(result["docx_style"]["output_path"]).paragraphs[1].style.name == "Heading 2"
    assert result["docx_break"]["object_type"] == "page_break"
    assert result["docx_table"]["object_type"] == "table"
    assert result["docx_merged"]["locator"].endswith(":row:0:cell:0")

    assert len(Presentation(result["pptx_added"]["output_path"]).slides) == 3
    assert result["pptx_duplicated"]["locator"] == "pptx:slide:2"
    assert Presentation(result["pptx_layout"]["output_path"]).slides[1].slide_layout.name == Presentation(result["pptx_layout"]["output_path"]).slide_layouts[5].name
    assert result["pptx_text"]["object_type"] == "text_shape"

    assert load_workbook(result["xlsx_write"]["output_path"])["Budget2026"]["B2"].value == 200
    assert result["xlsx_insert_rows"]["locator"] == "xlsx:sheet:Budget2026:row:2"
    assert result["xlsx_insert_columns"]["locator"] == "xlsx:sheet:Budget2026:col:1"
    assert result["xlsx_formula"]["locator"] == "xlsx:sheet:Budget2026:formula_cell:D4"
    assert "A4:B4" in {str(cell_range) for cell_range in load_workbook(result["xlsx_merged"]["output_path"])["Budget2026"].merged_cells.ranges}
