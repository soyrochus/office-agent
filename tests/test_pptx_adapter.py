from __future__ import annotations

from pptx import Presentation

from offagent.adapters import pptx_adapter


def test_extract_document_captures_text_shapes_and_metadata(sample_pptx) -> None:
    items = pptx_adapter.extract_document(sample_pptx)

    assert len(items) == 4
    assert all(item.item_id.startswith("slide:") for item in items)
    assert all(item.item_type == "slide_text_shape" for item in items)
    assert all("shape_id" in item.metadata for item in items)
    assert all("shape_index" in item.metadata for item in items)
    assert any(item.metadata["is_placeholder"] for item in items)

    summary_item = next(item for item in items if "Supplier shall present" in item.content_text)
    assert summary_item.content_text == "Alpha launch overview\nSupplier shall present the rollout plan."
    assert summary_item.metadata["slide_number"] == 1
    assert summary_item.metadata["text_frame_text"] == summary_item.content_text


def test_extract_document_excludes_non_text_shapes(sample_pptx) -> None:
    items = pptx_adapter.extract_document(sample_pptx)
    presentation = Presentation(str(sample_pptx))
    table_shape = next(shape for shape in presentation.slides[0].shapes if shape.has_table)
    table_item_id = pptx_adapter.make_item_id(1, table_shape.shape_id)

    assert table_item_id not in {item.item_id for item in items}


def test_build_embedding_text_returns_shape_content(sample_pptx) -> None:
    item = pptx_adapter.extract_document(sample_pptx)[0]

    assert pptx_adapter.build_embedding_text(item, sample_pptx) == item.content_text
