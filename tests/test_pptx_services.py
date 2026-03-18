from __future__ import annotations

import pytest
from pptx import Presentation

from offagent.adapters import pptx_adapter
from offagent.app.services import AppServices
from offagent.config import AppConfig


def test_index_search_locate_and_read_pptx(sample_pptx, tmp_path) -> None:
    services = AppServices(
        AppConfig(index_path=tmp_path / "state" / "index.sqlite3", document_roots=(tmp_path,))
    )

    summary = services.index_path(sample_pptx)
    hits = services.search_corpus("Supplier shall", file_type="pptx")
    slide_items = services.locate_slide_shapes(sample_pptx, 1)
    summary_item = next(item for item in slide_items if "Supplier shall present" in item.preview)
    read_text = services.read_item(sample_pptx, summary_item.item_id)

    assert summary.files_indexed == 1
    assert hits[0].item_id == summary_item.item_id
    assert hits[0].preview == "Alpha launch overview\nSupplier shall present the rollout plan."
    assert len(slide_items) == 3
    assert read_text == "Alpha launch overview\nSupplier shall present the rollout plan."


def test_replace_append_and_reindex_pptx(sample_pptx, tmp_path) -> None:
    services = AppServices(
        AppConfig(index_path=tmp_path / "state" / "index.sqlite3", document_roots=(tmp_path,))
    )
    services.index_document(sample_pptx)
    editable_item = services.locate_slide_shapes(sample_pptx, 2)[0]

    replace_result = services.replace_item_text(sample_pptx, editable_item.item_id, "Updated speaker notes.")
    append_result = services.append_item_text(sample_pptx, editable_item.item_id, "\nAdd follow-up action.")

    presentation = Presentation(str(sample_pptx))
    editable_shape = next(shape for shape in presentation.slides[1].shapes if shape.has_text_frame)

    assert replace_result.text == "Updated speaker notes."
    assert append_result.text == "Updated speaker notes.\nAdd follow-up action."
    assert editable_shape.text_frame.text == "Updated speaker notes.\nAdd follow-up action."

    hits = services.search_corpus("follow up action", file_type="pptx")
    assert hits[0].item_id == editable_item.item_id


def test_replace_rejects_non_text_pptx_target(sample_pptx, tmp_path) -> None:
    services = AppServices(
        AppConfig(index_path=tmp_path / "state" / "index.sqlite3", document_roots=(tmp_path,))
    )
    services.index_document(sample_pptx)

    presentation = Presentation(str(sample_pptx))
    table_shape = next(shape for shape in presentation.slides[0].shapes if shape.has_table)
    table_item_id = pptx_adapter.make_item_id(1, table_shape.shape_id)

    with pytest.raises(pptx_adapter.TargetNotEditableError, match="target not editable"):
        services.replace_item_text(sample_pptx, table_item_id, "Should fail")
