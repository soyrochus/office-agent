from __future__ import annotations

import asyncio
import sys
from pathlib import Path

import pytest
from docx import Document
from openpyxl import load_workbook
from openpyxl.cell.rich_text import CellRichText
from pptx import Presentation
from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client

from offagent.app.services import AppServices
from offagent.config import AppConfig
from offagent.domain.models import InlineFragment, InlineStyle, VisibleTextRange
from offagent.errors import TargetNotEditableError


def _services(tmp_path) -> AppServices:
    return AppServices(
        AppConfig(
            index_path=tmp_path / "state" / "index.sqlite3",
            document_roots=(tmp_path,),
            allowed_roots=(tmp_path,),
            output_roots=(tmp_path,),
        )
    )


def _server_params(config_path: Path) -> StdioServerParameters:
    return StdioServerParameters(
        command=sys.executable,
        args=["-m", "offagent", "mcp", "--config", str(config_path)],
    )


def _run_mcp(config_path: Path, callback):
    async def runner():
        async with stdio_client(_server_params(config_path)) as (read_stream, write_stream):
            async with ClientSession(read_stream, write_stream) as session:
                await session.initialize()
                return await callback(session)

    return asyncio.run(runner())


async def _call_tool(session: ClientSession, name: str, arguments: dict | None = None):
    result = await session.call_tool(name, arguments=arguments)
    assert not result.isError
    return result.structuredContent


def test_create_object_accepts_segments_for_docx_pptx_and_xlsx(sample_docx, sample_pptx, sample_xlsx, tmp_path) -> None:
    services = _services(tmp_path)
    docx_document = services.index_document(sample_docx)
    pptx_document = services.index_document(sample_pptx)
    xlsx_document = services.index_document(sample_xlsx)

    docx_created = services.create_object(
        docx_document.document_id,
        "docx:document",
        "paragraph",
        {},
        segments=[
            InlineFragment("Bold", InlineStyle(bold=True)),
            InlineFragment(" tail", InlineStyle()),
        ],
    )
    pptx_created = services.create_object(
        pptx_document.document_id,
        "pptx:slide:2",
        "text_shape",
        {"left": 914400, "top": 914400, "width": 1828800, "height": 914400},
        segments=[
            InlineFragment("Deck", InlineStyle(bold=True)),
            InlineFragment(" note", InlineStyle(italic=True)),
        ],
    )
    xlsx_created = services.create_object(
        xlsx_document.document_id,
        "xlsx:sheet:Notes 2026",
        "cell",
        {"coordinate": "B4"},
        segments=[
            InlineFragment("Rich", InlineStyle(bold=True)),
            InlineFragment(" text", InlineStyle(italic=True)),
        ],
    )

    updated_docx = Document(str(docx_created.output_path))
    assert [run.text for run in updated_docx.paragraphs[-1].runs] == ["Bold", " tail"]
    assert updated_docx.paragraphs[-1].runs[0].bold is True

    updated_pptx = Presentation(str(pptx_created.output_path))
    created_shape = next(shape for shape in updated_pptx.slides[1].shapes if shape.has_text_frame and shape.text == "Deck note")
    assert [run.text for run in created_shape.text_frame.paragraphs[0].runs] == ["Deck", " note"]
    assert created_shape.text_frame.paragraphs[0].runs[0].font.bold is True

    updated_xlsx = load_workbook(xlsx_created.output_path, rich_text=True)
    assert str(updated_xlsx["Notes 2026"]["B4"].value) == "Rich text"
    assert isinstance(updated_xlsx["Notes 2026"]["B4"].value, CellRichText)


def test_update_object_and_style_inline_support_partial_formatting_across_formats(
    sample_docx,
    sample_pptx,
    sample_xlsx,
    tmp_path,
) -> None:
    services = _services(tmp_path)
    docx_document = services.index_document(sample_docx)
    pptx_document = services.index_document(sample_pptx)
    xlsx_document = services.index_document(sample_xlsx)

    docx_updated = services.update_object(
        docx_document.document_id,
        "docx:para:1",
        {},
        segments=[
            InlineFragment("Alpha", InlineStyle(bold=True)),
            InlineFragment(" revised", InlineStyle(italic=True)),
        ],
    )
    docx_styled = services.style_inline(
        docx_updated.document_id,
        "docx:para:1",
        InlineStyle(underline=True),
        text_range=VisibleTextRange(start=6, end=13),
    )

    pptx_shape = next(
        child.locator
        for child in services.list_children(pptx_document.document_id, "pptx:slide:1")
        if child.object_type == "text_shape" and "Alpha launch overview" in child.preview
    )
    pptx_paragraph = services.list_children(
        pptx_document.document_id,
        pptx_shape,
        child_type="paragraph",
    )[1].locator
    pptx_updated = services.update_object(
        pptx_document.document_id,
        pptx_paragraph,
        {},
        segments=[
            InlineFragment("Supplier", InlineStyle(bold=True)),
            InlineFragment(" update", InlineStyle(italic=True)),
        ],
    )
    pptx_styled = services.style_inline(
        pptx_updated.document_id,
        pptx_paragraph,
        InlineStyle(font_color="AA2200"),
        text_range=VisibleTextRange(start=0, end=8),
    )

    xlsx_updated = services.update_object(
        xlsx_document.document_id,
        "xlsx:sheet:Notes 2026!A1",
        {},
        segments=[
            InlineFragment("Supplier", InlineStyle(bold=True)),
            InlineFragment(" shall review variance.", InlineStyle()),
        ],
    )
    xlsx_styled = services.style_inline(
        xlsx_updated.document_id,
        "xlsx:sheet:Notes 2026!A1",
        InlineStyle(italic=True),
        text_range=VisibleTextRange(start=9, end=14),
    )

    rewritten_docx = Document(str(docx_styled.output_path))
    assert rewritten_docx.paragraphs[1].text == "Alpha revised"
    assert [run.text for run in rewritten_docx.paragraphs[1].runs] == ["Alpha", " ", "revised"]
    assert rewritten_docx.paragraphs[1].runs[2].underline is True

    rewritten_pptx = Presentation(str(pptx_styled.output_path))
    target_shape = next(shape for shape in rewritten_pptx.slides[0].shapes if shape.has_text_frame and "Supplier update" in shape.text)
    paragraph = target_shape.text_frame.paragraphs[1]
    assert paragraph.text == "Supplier update"
    assert paragraph.runs[0].font.color.rgb is not None

    rewritten_xlsx = load_workbook(xlsx_styled.output_path, rich_text=True)
    cell = rewritten_xlsx["Notes 2026"]["A1"]
    assert str(cell.value) == "Supplier shall review variance."
    assert isinstance(cell.value, CellRichText)


def test_partial_formatting_rejects_unsupported_targets(sample_pptx, sample_xlsx, tmp_path) -> None:
    services = _services(tmp_path)
    pptx_document = services.index_document(sample_pptx)
    xlsx_document = services.index_document(sample_xlsx)

    table_locator = next(
        child.locator
        for child in services.list_children(pptx_document.document_id, "pptx:slide:1")
        if child.object_type == "table"
    )

    with pytest.raises(TargetNotEditableError, match="editable"):
        services.style_inline(
            pptx_document.document_id,
            table_locator,
            InlineStyle(bold=True),
            text_range=VisibleTextRange(start=0, end=5),
        )

    with pytest.raises(TargetNotEditableError, match="non-string|formula"):
        services.style_inline(
            xlsx_document.document_id,
            "xlsx:sheet:Budget2026!B2",
            InlineStyle(bold=True),
            text_range=VisibleTextRange(start=0, end=3),
        )


def test_mcp_exposes_partial_formatting_schema_and_round_trips(sample_docx, sample_xlsx, config_path) -> None:
    def scenario(session: ClientSession):
        async def run():
            tools = await session.list_tools()
            tool_map = {tool.name: tool for tool in tools.tools}
            assert "segments" in tool_map["create_object"].inputSchema["properties"]
            assert "segments" in tool_map["update_object"].inputSchema["properties"]
            assert "range" in tool_map["style_inline"].inputSchema["properties"]

            await _call_tool(session, "index_documents", {"paths": [str(sample_docx), str(sample_xlsx)]})
            documents = await _call_tool(session, "list_documents")
            docs_by_type = {document["file_type"]: document for document in documents["documents"]}

            updated = await _call_tool(
                session,
                "update_object",
                {
                    "document_id": docs_by_type["docx"]["document_id"],
                    "locator": "docx:para:1",
                    "properties": {},
                    "segments": [
                        {"text": "Alpha", "style": {"bold": True}},
                        {"text": " MCP", "style": {"italic": True}},
                    ],
                },
            )
            styled = await _call_tool(
                session,
                "style_inline",
                {
                    "document_id": updated["document_id"],
                    "locator": "docx:para:1",
                    "style": {"underline": True},
                    "range": {"start": 6, "end": 9},
                },
            )
            failure = await session.call_tool(
                "style_inline",
                arguments={
                    "document_id": docs_by_type["xlsx"]["document_id"],
                    "locator": "xlsx:sheet:Budget2026!B2",
                    "style": {"bold": True},
                    "range": {"start": 0, "end": 3},
                },
            )
            assert failure.isError
            return styled

        return run()

    result = _run_mcp(config_path, scenario)
    assert result["locator"] == "docx:para:1"
