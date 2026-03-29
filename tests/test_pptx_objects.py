from __future__ import annotations

import base64
from pathlib import Path

import pytest
from pptx import Presentation
from pptx.util import Inches

from offagent.domain.models import Capability
from offagent.errors import InvalidArgumentsError
from offagent.objects.pptx_objects import (
    PptxObjectResolver,
    add_slide,
    add_text_shape,
    duplicate_slide,
    set_slide_layout,
)

PNG_1X1 = (
    "iVBORw0KGgoAAAANSUhEUgAAAAEAAAABCAQAAAC1HAwCAAAAC0lEQVR42mP8/x8AAwMCAO+/aRsAAAAASUVORK5CYII="
)


@pytest.fixture
def pptx_with_image(tmp_path) -> Path:
    image_path = tmp_path / "dot.png"
    image_path.write_bytes(base64.b64decode(PNG_1X1))

    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    slide.shapes.add_textbox(Inches(1), Inches(1), Inches(4), Inches(1)).text = "Image slide"
    slide.shapes.add_picture(str(image_path), Inches(1), Inches(2), Inches(1), Inches(1))

    path = tmp_path / "image.pptx"
    presentation.save(path)
    return path


def test_get_presentation_slide_and_notes_objects(sample_pptx) -> None:
    resolver = PptxObjectResolver()

    presentation = resolver.get_object(sample_pptx, "pptx:presentation")
    slide = resolver.get_object(sample_pptx, "pptx:slide:1")
    notes = resolver.get_object(sample_pptx, "pptx:slide:1:notes")

    assert presentation.object_type == "presentation"
    assert presentation.properties["slide_count"] == 2
    assert slide.object_type == "slide"
    assert slide.properties["slide_number"] == 1
    assert notes.object_type == "notes"
    assert "Speaker notes" in notes.properties["text"]


def test_get_text_shape_table_row_and_cell_objects(sample_pptx) -> None:
    resolver = PptxObjectResolver()

    slide_children = resolver.list_children(sample_pptx, "pptx:slide:1")
    text_shape_locator = next(child.locator for child in slide_children if child.object_type == "text_shape")
    table_locator = next(child.locator for child in slide_children if child.object_type == "table")

    text_shape = resolver.get_object(sample_pptx, text_shape_locator)
    table = resolver.get_object(sample_pptx, table_locator)
    row = resolver.get_object(sample_pptx, f"{table_locator}:row:0")
    cell = resolver.get_object(sample_pptx, f"{table_locator}:row:0:cell:0")

    assert text_shape.object_type == "text_shape"
    assert "text" in text_shape.properties
    assert table.object_type == "table"
    assert table.child_summary[0].locator == f"{table_locator}:row:0"
    assert row.object_type == "table_row"
    assert cell.object_type == "table_cell"
    assert cell.properties["text"] == "Table text to ignore"


def test_list_children_filters_slide_and_table_objects(sample_pptx) -> None:
    resolver = PptxObjectResolver()

    slides = resolver.list_children(sample_pptx, "pptx:presentation", child_type="slide")
    slide_shapes = resolver.list_children(sample_pptx, "pptx:slide:1", child_type="text_shape")
    table_locator = next(child.locator for child in resolver.list_children(sample_pptx, "pptx:slide:1") if child.object_type == "table")
    rows = resolver.list_children(sample_pptx, table_locator, child_type="table_row")

    assert [child.locator for child in slides] == ["pptx:slide:1", "pptx:slide:2"]
    assert slide_shapes
    assert all(child.object_type == "text_shape" for child in slide_shapes)
    assert len(rows) == 2


def test_resolve_capabilities_for_slide_matches_spec(sample_pptx) -> None:
    resolver = PptxObjectResolver()

    capabilities = resolver.resolve_capabilities(sample_pptx, "pptx:slide:1")

    assert capabilities == frozenset(
        {
            Capability.READ,
            Capability.UPDATE,
            Capability.DELETE,
            Capability.ADD_CHILD,
            Capability.MOVE,
            Capability.COPY,
        }
    )


def test_get_image_shape_object(pptx_with_image) -> None:
    resolver = PptxObjectResolver()

    image_shape_locator = next(
        child.locator
        for child in resolver.list_children(pptx_with_image, "pptx:slide:1")
        if child.object_type == "image_shape"
    )
    image_shape = resolver.get_object(pptx_with_image, image_shape_locator)

    assert image_shape.object_type == "image_shape"
    assert image_shape.properties["shape_id"] > 0


def test_pptx_escape_hatches_add_duplicate_relayout_and_add_text_shape(sample_pptx, tmp_path) -> None:
    added_path = tmp_path / "added.pptx"
    duplicated_path = tmp_path / "duplicated.pptx"
    layout_path = tmp_path / "layout.pptx"
    text_path = tmp_path / "text-shape.pptx"
    resolver = PptxObjectResolver()

    added_locator, _, _ = add_slide(sample_pptx, layout_index=6, output_path=added_path)
    duplicated_locator, _, _ = duplicate_slide(
        sample_pptx,
        "pptx:slide:1",
        position=2,
        output_path=duplicated_path,
    )
    relayout_locator, _, _ = set_slide_layout(
        sample_pptx,
        "pptx:slide:2",
        layout_index=5,
        output_path=layout_path,
    )
    text_locator, _, _ = add_text_shape(
        sample_pptx,
        "pptx:slide:2",
        text="Inserted textbox",
        left=914400,
        top=914400,
        width=1828800,
        height=914400,
        output_path=text_path,
    )

    assert added_locator == "pptx:slide:3"
    assert len(Presentation(str(added_path)).slides) == 3

    duplicated = Presentation(str(duplicated_path))
    assert duplicated_locator == "pptx:slide:2"
    assert len(duplicated.slides) == 3
    assert any("Quarterly Planning" in shape.text for shape in duplicated.slides[1].shapes if shape.has_text_frame)

    relayout = Presentation(str(layout_path))
    assert relayout_locator == "pptx:slide:2"
    assert relayout.slides[1].slide_layout.name == relayout.slide_layouts[5].name

    added_text = resolver.get_object(text_path, text_locator)
    assert added_text.object_type == "text_shape"
    assert added_text.properties["text"] == "Inserted textbox"


def test_pptx_add_slide_rejects_unknown_layout(sample_pptx, tmp_path) -> None:
    with pytest.raises(InvalidArgumentsError, match="Unknown PPTX slide layout index"):
        add_slide(sample_pptx, layout_index=99, output_path=tmp_path / "invalid-layout.pptx")
