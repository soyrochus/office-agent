from __future__ import annotations

import shutil
from pathlib import Path

import pytest
from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pptx.util import Inches

FIXTURES_DIR = Path(__file__).parent / "fixtures"


@pytest.fixture
def sample_docx(tmp_path) -> Path:
    path = tmp_path / "sample.docx"
    document = Document()
    document.add_heading("Project Heading", level=1)
    paragraph = document.add_paragraph()
    first_run = paragraph.add_run("Alpha paragraph for search.")
    first_run.bold = True
    document.add_paragraph("")
    document.add_paragraph("Supplier shall deliver by Friday.")
    table = document.add_table(rows=1, cols=1)
    table.cell(0, 0).text = "Table text to ignore"
    document.save(path)
    return path


@pytest.fixture
def config_path(tmp_path) -> Path:
    path = tmp_path / "office-agent.toml"
    path.write_text(
        f"""
[offagent]
index_path = "{(tmp_path / 'state' / 'index.sqlite3').as_posix()}"
document_roots = ["{tmp_path.as_posix()}"]
embedding_model = "hash://test"
embedding_dimensions = 48
""".strip()
    )
    return path


@pytest.fixture
def sample_pptx(tmp_path) -> Path:
    path = tmp_path / "sample.pptx"
    presentation = Presentation()

    first_slide = presentation.slides.add_slide(presentation.slide_layouts[5])
    title = first_slide.shapes.title
    title.text = "Quarterly Planning"

    summary_box = first_slide.shapes.add_textbox(Inches(1), Inches(1.8), Inches(5), Inches(1.5))
    summary_frame = summary_box.text_frame
    summary_frame.text = "Alpha launch overview"
    summary_frame.add_paragraph().text = "Supplier shall present the rollout plan."

    notes_box = first_slide.shapes.add_textbox(Inches(1), Inches(3.7), Inches(4.5), Inches(1))
    notes_box.text = "Secondary text frame"
    first_slide.notes_slide.notes_text_frame.text = "Speaker notes: confirm launch owner."

    table_shape = first_slide.shapes.add_table(2, 2, Inches(6), Inches(1.8), Inches(3), Inches(1.5))
    table_shape.table.cell(0, 0).text = "Table text to ignore"

    second_slide = presentation.slides.add_slide(presentation.slide_layouts[6])
    editable_box = second_slide.shapes.add_textbox(Inches(1), Inches(1.5), Inches(5.5), Inches(1.5))
    editable_box.text = "Editable speaker notes"
    second_slide.notes_slide.notes_text_frame.text = "Speaker notes: review action list."

    presentation.save(path)
    return path


@pytest.fixture
def sample_xlsx(tmp_path) -> Path:
    path = tmp_path / "sample.xlsx"
    workbook = Workbook()

    budget_sheet = workbook.active
    budget_sheet.title = "Budget2026"
    budget_sheet["A1"] = "Quarterly Budget"
    budget_sheet["B2"] = 125000
    budget_sheet["C3"] = "=SUM(1,2)"

    notes_sheet = workbook.create_sheet("Notes 2026")
    notes_sheet["A1"] = "Supplier shall review variance."
    notes_sheet["A2"] = "Follow up with finance."

    workbook.save(path)
    return path


def _copy_fixture(source_name: str, target_dir: Path) -> Path:
    source = FIXTURES_DIR / source_name
    target = target_dir / source_name
    shutil.copy2(source, target)
    return target


@pytest.fixture
def golden_workspace(tmp_path) -> Path:
    workspace = tmp_path / "workspace"
    docs_dir = workspace / "docs"
    docs_dir.mkdir(parents=True)
    _copy_fixture("golden.docx", docs_dir)
    _copy_fixture("golden.pptx", docs_dir)
    _copy_fixture("golden.xlsx", docs_dir)
    return workspace


@pytest.fixture
def golden_docx(golden_workspace) -> Path:
    return golden_workspace / "docs" / "golden.docx"


@pytest.fixture
def golden_pptx(golden_workspace) -> Path:
    return golden_workspace / "docs" / "golden.pptx"


@pytest.fixture
def golden_xlsx(golden_workspace) -> Path:
    return golden_workspace / "docs" / "golden.xlsx"


@pytest.fixture
def golden_config_path(golden_workspace) -> Path:
    path = golden_workspace / "office-agent.toml"
    docs_dir = golden_workspace / "docs"
    output_dir = golden_workspace / "edited"
    path.write_text(
        f"""
[offagent]
index_path = "{(golden_workspace / 'state' / 'index.sqlite3').as_posix()}"
document_roots = ["{docs_dir.as_posix()}"]
allowed_roots = ["{docs_dir.as_posix()}"]
output_directory = "{output_dir.as_posix()}"
output_roots = ["{output_dir.as_posix()}", "{docs_dir.as_posix()}"]
embedding_model = "hash://golden"
embedding_dimensions = 48
""".strip()
    )
    return path
