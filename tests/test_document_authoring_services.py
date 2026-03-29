from __future__ import annotations

import pytest
from docx import Document
from openpyxl import load_workbook
from pptx import Presentation

from offagent.app.services import AppServices
from offagent.config import AppConfig
from offagent.domain.models import BlockStyle, InlineStyle
from offagent.errors import InvalidArgumentsError


def _services(tmp_path) -> AppServices:
    return AppServices(
        AppConfig(
            index_path=tmp_path / "state" / "index.sqlite3",
            document_roots=(tmp_path,),
            allowed_roots=(tmp_path,),
            output_roots=(tmp_path,),
        )
    )


def test_create_document_indexes_formats_and_seeds_defaults(tmp_path) -> None:
    services = _services(tmp_path)

    docx_result = services.create_document("docx", tmp_path / "report.docx")
    pptx_result = services.create_document("pptx", tmp_path / "deck.pptx")
    xlsx_result = services.create_document("xlsx", tmp_path / "sheet.xlsx")
    named_xlsx_result = services.create_document(
        "xlsx",
        tmp_path / "named.xlsx",
        initial_sheet_name="Budget",
    )

    assert docx_result.output_path is not None
    assert pptx_result.output_path is not None
    assert xlsx_result.output_path is not None
    assert named_xlsx_result.output_path is not None
    assert ".edited." in docx_result.output_path.name
    assert docx_result.output_path.suffix == ".docx"
    assert pptx_result.output_path.suffix == ".pptx"
    assert xlsx_result.output_path.suffix == ".xlsx"

    docx_document = Document(str(docx_result.output_path))
    style_names = {style.name for style in docx_document.styles if getattr(style, "name", None)}
    assert {"Normal", "Heading 1", "Title", "Caption"} <= style_names
    assert services.get_document(docx_result.document_id).path == docx_result.output_path

    presentation = Presentation(str(pptx_result.output_path))
    assert len(presentation.slide_layouts) >= 1
    assert services.get_document(pptx_result.document_id).path == pptx_result.output_path

    workbook = load_workbook(xlsx_result.output_path)
    named_workbook = load_workbook(named_xlsx_result.output_path)
    assert workbook.sheetnames == ["Sheet1"]
    assert named_workbook.sheetnames == ["Budget"]
    assert services.get_document(xlsx_result.document_id).path == xlsx_result.output_path


def test_add_content_block_supports_all_document_pairs_and_rejects_invalid_combo(tmp_path) -> None:
    services = _services(tmp_path)

    docx = services.create_document("docx", tmp_path / "author.docx")
    paragraph = services.add_content_block(
        docx.document_id,
        "paragraph",
        {"text": "Alpha paragraph"},
    )
    heading = services.add_content_block(
        paragraph.document_id,
        "heading",
        {"text": "Section title", "level": 2},
    )
    table = services.add_content_block(
        heading.document_id,
        "table",
        {"rows": 2, "columns": 3},
    )

    docx_document = Document(str(table.output_path))
    assert [paragraph.text for paragraph in docx_document.paragraphs[:2]] == [
        "Alpha paragraph",
        "Section title",
    ]
    assert docx_document.paragraphs[1].style.name == "Heading 2"
    assert len(docx_document.tables) == 1
    assert table.locator == "docx:table:0"

    pptx = services.create_document("pptx", tmp_path / "author.pptx")
    slide = services.add_content_block(pptx.document_id, "slide", {})
    textbox = services.add_content_block(
        slide.document_id,
        "textbox",
        {"slide_locator": slide.locator, "text": "Presenter notes"},
    )

    presentation = Presentation(str(textbox.output_path))
    text_shapes = [shape for shape in presentation.slides[0].shapes if shape.has_text_frame]
    assert len(presentation.slides) == 1
    assert any(shape.text_frame.text == "Presenter notes" for shape in text_shapes)
    assert textbox.locator is not None and textbox.locator.startswith("pptx:slide:1:shape:")

    xlsx = services.create_document("xlsx", tmp_path / "author.xlsx")
    sheet = services.add_content_block(xlsx.document_id, "sheet", {"name": "Data"})
    row = services.add_content_block(
        sheet.document_id,
        "row",
        {"sheet_locator": "xlsx:sheet:Data", "values": ["A", "B", "C"]},
    )
    cell = services.add_content_block(
        row.document_id,
        "cell",
        {"cell_locator": "xlsx:sheet:Data!B3", "value": "Written"},
    )

    workbook = load_workbook(cell.output_path)
    assert workbook.sheetnames == ["Sheet1", "Data"]
    assert [workbook["Data"][coordinate].value for coordinate in ("A1", "B1", "C1")] == ["A", "B", "C"]
    assert workbook["Data"]["B3"].value == "Written"
    assert row.locator == "xlsx:sheet:Data:row:1"
    assert cell.locator == "xlsx:sheet:Data!B3"

    with pytest.raises(InvalidArgumentsError, match="Unsupported add_content_block combination"):
        services.add_content_block(slide.document_id, "paragraph", {"text": "bad"})


def test_style_inline_applies_docx_pptx_xlsx_and_supports_clear_fields(
    sample_docx,
    sample_pptx,
    sample_xlsx,
    tmp_path,
) -> None:
    services = _services(tmp_path)
    docx_ref = services.index_document(sample_docx)
    pptx_ref = services.index_document(sample_pptx)
    xlsx_ref = services.index_document(sample_xlsx)

    docx_result = services.style_inline(
        docx_ref.document_id,
        "docx:para:1:run:0",
        InlineStyle(
            bold=False,
            italic=True,
            font_name="Calibri",
            font_color="00AA00",
        ),
        clear_fields=["bold"],
    )
    docx_document = Document(str(docx_result.output_path))
    docx_run = docx_document.paragraphs[1].runs[0]
    assert docx_run.bold is None
    assert docx_run.italic is True
    assert docx_run.font.name == "Calibri"
    assert str(docx_run.font.color.rgb) == "00AA00"

    shape_locator = services.locate_slide_shapes(sample_pptx, 2)[0].item_id
    pptx_result = services.style_inline(
        pptx_ref.document_id,
        shape_locator,
        InlineStyle(bold=True, font_name="Aptos", font_size=18, font_color="112233"),
    )
    presentation = Presentation(str(pptx_result.output_path))
    editable_shape = next(shape for shape in presentation.slides[1].shapes if shape.has_text_frame)
    run = editable_shape.text_frame.paragraphs[0].runs[0]
    assert run.font.bold is True
    assert run.font.name == "Aptos"
    assert round(run.font.size.pt) == 18
    assert str(run.font.color.rgb) == "112233"

    xlsx_result = services.style_inline(
        xlsx_ref.document_id,
        "xlsx:sheet:Notes 2026!A1",
        InlineStyle(bold=True, italic=True, font_name="Calibri", font_color="334455"),
    )
    workbook = load_workbook(xlsx_result.output_path)
    cell = workbook["Notes 2026"]["A1"]
    assert cell.font.bold is True
    assert cell.font.italic is True
    assert cell.font.name == "Calibri"
    assert cell.font.color.rgb.endswith("334455")


def test_style_block_applies_docx_pptx_xlsx_properties(sample_docx, sample_pptx, sample_xlsx, tmp_path) -> None:
    services = _services(tmp_path)
    docx_ref = services.index_document(sample_docx)
    pptx_ref = services.index_document(sample_pptx)
    xlsx_ref = services.index_document(sample_xlsx)

    docx_result = services.style_block(
        docx_ref.document_id,
        "docx:para:1",
        BlockStyle(alignment="center", spacing_before=12, spacing_after=6),
    )
    docx_document = Document(str(docx_result.output_path))
    paragraph = docx_document.paragraphs[1]
    assert paragraph.alignment is not None
    assert round(paragraph.paragraph_format.space_before.pt) == 12
    assert round(paragraph.paragraph_format.space_after.pt) == 6

    shape_locator = services.locate_slide_shapes(sample_pptx, 2)[0].item_id
    pptx_result = services.style_block(
        pptx_ref.document_id,
        shape_locator,
        BlockStyle(alignment="center", indent_level=1),
    )
    presentation = Presentation(str(pptx_result.output_path))
    editable_shape = next(shape for shape in presentation.slides[1].shapes if shape.has_text_frame)
    pptx_paragraph = editable_shape.text_frame.paragraphs[0]
    assert pptx_paragraph.alignment is not None
    assert pptx_paragraph.level == 1

    xlsx_result = services.style_block(
        xlsx_ref.document_id,
        "xlsx:sheet:Notes 2026!A1",
        BlockStyle(alignment="center", wrap_text=True, fill_color="FFEE00"),
    )
    workbook = load_workbook(xlsx_result.output_path)
    cell = workbook["Notes 2026"]["A1"]
    assert cell.alignment.horizontal == "center"
    assert cell.alignment.wrap_text is True
    assert cell.fill.fill_type == "solid"
    assert cell.fill.start_color.rgb.endswith("FFEE00")


@pytest.mark.parametrize(
    ("role", "level", "expected_style"),
    [
        ("title", None, "Title"),
        ("body", None, "Normal"),
        ("caption", None, "Caption"),
        ("heading", 3, "Heading 3"),
    ],
)
def test_set_structural_role_maps_roles_and_requires_docx(
    role: str,
    level: int | None,
    expected_style: str,
    tmp_path,
) -> None:
    services = _services(tmp_path)
    docx = services.create_document("docx", tmp_path / f"{role}.docx")
    paragraph = services.add_content_block(docx.document_id, "paragraph", {"text": "Styled"})

    result = services.set_structural_role(paragraph.document_id, "docx:para:0", role, level=level)
    document = Document(str(result.output_path))
    assert document.paragraphs[0].style.name == expected_style


def test_set_structural_role_rejects_invalid_heading_level_and_non_docx(tmp_path) -> None:
    services = _services(tmp_path)
    docx = services.create_document("docx", tmp_path / "reject.docx")
    paragraph = services.add_content_block(docx.document_id, "paragraph", {"text": "Styled"})
    pptx = services.create_document("pptx", tmp_path / "reject.pptx")
    xlsx = services.create_document("xlsx", tmp_path / "reject.xlsx")

    with pytest.raises(InvalidArgumentsError, match="Heading structural role requires level"):
        services.set_structural_role(paragraph.document_id, "docx:para:0", "heading")

    with pytest.raises(InvalidArgumentsError, match="requires a .docx document"):
        services.set_structural_role(pptx.document_id, "pptx:slide:1", "title")

    with pytest.raises(InvalidArgumentsError, match="requires a .docx document"):
        services.set_structural_role(xlsx.document_id, "xlsx:sheet:Sheet1!A1", "title")


def test_new_authoring_operations_use_versioned_outputs_and_reindex(tmp_path) -> None:
    services = _services(tmp_path)

    created = services.create_document("docx", tmp_path / "flow.docx")
    paragraph = services.add_content_block(created.document_id, "paragraph", {"text": "Flow text"})
    inline = services.style_inline(
        paragraph.document_id,
        "docx:para:0:run:0",
        InlineStyle(italic=True),
    )
    block = services.style_block(
        inline.document_id,
        "docx:para:0",
        BlockStyle(alignment="center"),
    )
    role = services.set_structural_role(block.document_id, "docx:para:0", "heading", level=2)

    assert created.output_path is not None and ".edited." in created.output_path.name
    assert paragraph.output_path is not None and ".edited." in paragraph.output_path.name
    assert inline.output_path is not None and ".edited." in inline.output_path.name
    assert block.output_path is not None and ".edited." in block.output_path.name
    assert role.output_path is not None and ".edited." in role.output_path.name

    final_document = services.get_document(role.document_id)
    assert final_document.path == role.output_path
    final_object = services.get_object(role.document_id, "docx:para:0")
    assert final_object.properties["style_name"] == "Heading 2"
